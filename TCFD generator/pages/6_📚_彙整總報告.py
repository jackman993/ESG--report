"""Step 4: 彙整總報告"""
import streamlit as st
import sys
import json
import io
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree

# 導入共享模組
sys.path.insert(0, str(Path(__file__).parent.parent))
from shared.config import *
from shared.utils import render_output_folder_links, render_api_key_input, render_sidebar_navigation

# ============ PPTX 合併函數 ============
def find_latest_pptx(directory):
    """找到目錄中最新的 PPTX 文件"""
    if not directory.exists():
        return None
    
    pptx_files = list(directory.glob("*.pptx"))
    if not pptx_files:
        return None
    
    # 按修改時間排序，返回最新的
    latest = max(pptx_files, key=lambda p: p.stat().st_mtime)
    return latest

def normalize_fonts_in_slide(slide, target_font="Microsoft JhengHei"):
    """統一投影片中的字體（避免字體不一致導致修復提示）"""
    try:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font and run.font.name:
                            # 統一替換字體
                            run.font.name = target_font
    except Exception as e:
        # 字體統一失敗不影響合併
        pass

def merge_pptx_files(file_paths, output_path):
    """合併多個 PPTX 文件（統一處理字體，避免修復提示）"""
    if not file_paths:
        raise ValueError("沒有文件可以合併")
    
    # 使用第一個文件作為基礎
    base_prs = Presentation(str(file_paths[0]))
    
    # 刪除第一個文件的所有投影片（我們要重新添加）
    while len(base_prs.slides) > 0:
        rId = base_prs.slides._sldIdLst[0].rId
        base_prs.part.drop_rel(rId)
        del base_prs.slides._sldIdLst[0]
    
    total_slides = 0
    
    # 統一使用的字體（使用環境段的字體，因為它最完整）
    unified_font = "Microsoft JhengHei"
    
    for file_path in file_paths:
        if not file_path or not file_path.exists():
            st.warning(f"⚠️ 跳過不存在的文件：{file_path}")
            continue
        
        try:
            source_prs = Presentation(str(file_path))
            
            # 確保簡報尺寸一致
            if total_slides == 0:
                base_prs.slide_width = source_prs.slide_width
                base_prs.slide_height = source_prs.slide_height
            
            # 使用更可靠的方法：直接複製投影片的完整 XML
            # 安全地獲取空白版面（避免索引超出範圍）
            blank_layout = None
            
            # 首先檢查是否有可用的布局
            if len(base_prs.slide_layouts) == 0:
                st.error(f"❌ 模板文件沒有可用的版面配置：{file_path.name}")
                continue
            
            # 嘗試找到空白布局
            for layout in base_prs.slide_layouts:
                try:
                    name = (layout.name or "").lower()
                except Exception:
                    name = ""
                if "blank" in name or "空白" in name or "title only" in name:
                    blank_layout = layout
                    break
            
            # 如果找不到空白布局，使用第一個可用的布局
            if blank_layout is None:
                try:
                    blank_layout = base_prs.slide_layouts[0]
                except IndexError:
                    st.error(f"❌ 無法獲取版面配置：{file_path.name}")
                    continue
            
            for slide in source_prs.slides:
                # 創建新的投影片（使用空白版面）
                new_slide = base_prs.slides.add_slide(blank_layout)
                
                # 獲取原始投影片的 XML
                source_xml = slide.element
                
                # 清空新投影片的預設內容
                for shape in list(new_slide.shapes):
                    sp = shape._element
                    sp.getparent().remove(sp)
                
                # 直接複製整個投影片的 XML（包括版面配置）
                # 複製 cSld (common slide data) - 這是投影片的主要內容
                source_cSld = source_xml.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
                if source_cSld is not None:
                    # 找到新投影片的 cSld
                    new_cSld = new_slide.element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
                    if new_cSld is not None:
                        # 清空新投影片的 cSld
                        for child in list(new_cSld):
                            new_cSld.remove(child)
                        # 複製所有子元素（深層複製）
                        for child in source_cSld:
                            new_cSld.append(etree.fromstring(etree.tostring(child, encoding='unicode').encode('utf-8')))
                
                # 複製投影片的關係（圖片、媒體等）- 這很重要！
                # 需要先複製關係，再複製 XML，這樣圖片引用才能正確
                rel_map = {}  # 映射原始關係 ID 到新關係 ID
                
                for rel in slide.part.rels.values():
                    try:
                        if rel.is_external:
                            continue
                        
                        # 獲取關係的目標部分（圖片數據）
                        target_part = rel.target_part
                        
                        # 複製關係到新投影片
                        new_rel = new_slide.part.rels.add_relationship(
                            rel.rtype,
                            rel.target_ref,
                            target_part
                        )
                        
                        # 記錄關係映射（原始 ID -> 新 ID）
                        rel_map[rel.rId] = new_rel.rId
                        
                    except Exception as rel_error:
                        # 如果關係複製失敗，嘗試直接複製圖片 blob
                        try:
                            if 'image' in rel.target_ref.lower() or rel.rtype.endswith('image'):
                                # 嘗試從形狀中直接複製圖片
                                for shape in slide.shapes:
                                    if hasattr(shape, 'image') and shape.image:
                                        try:
                                            # 直接複製圖片到新投影片
                                            left = shape.left
                                            top = shape.top
                                            width = shape.width
                                            height = shape.height
                                            image_blob = shape.image.blob
                                            new_slide.shapes.add_picture(
                                                io.BytesIO(image_blob),
                                                left, top, width, height
                                            )
                                        except:
                                            pass
                        except:
                            pass
                
                # 更新 XML 中的關係引用（將原始關係 ID 替換為新的關係 ID）
                if rel_map:
                    for old_rId, new_rId in rel_map.items():
                        # 在投影片 XML 中替換關係引用
                        xml_str = etree.tostring(new_slide.element, encoding='unicode')
                        xml_str = xml_str.replace(f'rId="{old_rId}"', f'rId="{new_rId}"')
                        xml_str = xml_str.replace(f'r:id="{old_rId}"', f'r:id="{new_rId}"')
                        new_slide.element = etree.fromstring(xml_str.encode('utf-8'))
                
                # 統一字體（避免字體不一致導致修復提示）
                normalize_fonts_in_slide(new_slide, unified_font)
                
                total_slides += 1
            
            st.success(f"✅ 已合併：{file_path.name} ({len(source_prs.slides)} 頁，字體已統一為 {unified_font})")
            
        except Exception as e:
            st.error(f"❌ 合併 {file_path.name} 時出錯：{e}")
            import traceback
            st.code(traceback.format_exc())
            # 繼續處理下一個文件
            continue
    
    # 儲存合併後的簡報
    base_prs.save(str(output_path))
    return total_slides

# 頁面配置
st.set_page_config(page_title="Step 4: 彙整總報告", page_icon="📚", layout="wide")

# 側邊欄（自定義導航）
render_sidebar_navigation()
st.sidebar.divider()
API_KEY = render_api_key_input()
render_output_folder_links()

# 主頁面
st.title("📚 Step 4: 彙整總報告")

# 前置條件檢查
st.subheader("📋 前置條件檢查")

# 檢查1：從 session_state 讀取標誌
step1_done_flag = st.session_state.get("step1_done", False)
step2_done_flag = st.session_state.get("step2_done", False)
step3_done_flag = st.session_state.get("step3_done", False)

# 檢查2：檢查實際文件是否存在（備用驗證，確保可靠性）
env_file = find_latest_pptx(OUTPUT_C_ENVIRONMENT)
company_file = find_latest_pptx(OUTPUT_D_COMPANY)
govsoci_file = find_latest_pptx(OUTPUT_F_GOVSOCI)

# 最終判斷：標誌為True 或 文件存在
step1_done = step1_done_flag or (env_file is not None and env_file.exists())
step2_done = step2_done_flag or (company_file is not None and company_file.exists())
step3_done = step3_done_flag or (govsoci_file is not None and govsoci_file.exists())

# 如果文件存在但標誌未設置，自動恢復標誌
if step1_done and not step1_done_flag:
    st.session_state.step1_done = True
if step2_done and not step2_done_flag:
    st.session_state.step2_done = True
if step3_done and not step3_done_flag:
    st.session_state.step3_done = True

col1, col2, col3 = st.columns(3)
with col1:
    st.success("✅ Step 1") if step1_done else st.warning("⬜ Step 1")
with col2:
    st.success("✅ Step 2") if step2_done else st.warning("⬜ Step 2")
with col3:
    st.success("✅ Step 3") if step3_done else st.warning("⬜ Step 3")

st.divider()

# 報告合併
if all([step1_done, step2_done, step3_done]):
    st.subheader("📑 報告合併")
    
    st.info("""
    **合併順序：**
    1. 重大議題與公司段（Step 2）
    2. 環境段（Step 1 子步驟3 的輸出）
    3. 治理與社會段（Step 3）
    """)
    
    if st.button("🚀 彙整總報告", type="primary", use_container_width=True):
        with st.spinner("📄 正在合併所有報告..."):
            try:
                # 找到最新的文件
                company_file = find_latest_pptx(OUTPUT_D_COMPANY)
                environment_file = find_latest_pptx(OUTPUT_C_ENVIRONMENT)
                govsoci_file = find_latest_pptx(OUTPUT_F_GOVSOCI)
                
                # 檢查文件是否存在
                files_to_merge = []
                file_names = []
                
                # 合併順序：1. 重大議題與公司段（Step 2） 2. 環境段（Step 1） 3. 治理與社會段（Step 3）
                if company_file:
                    files_to_merge.append(company_file)
                    file_names.append(f"重大議題與公司段（Step 2）：{company_file.name}")
                else:
                    st.warning("⚠️ 未找到重大議題與公司段報告（Step 2）")
                
                if environment_file:
                    files_to_merge.append(environment_file)
                    file_names.append(f"環境段（Step 1 子步驟3）：{environment_file.name}")
                else:
                    st.warning("⚠️ 未找到環境段報告（Step 1 子步驟3）")
                
                if govsoci_file:
                    files_to_merge.append(govsoci_file)
                    file_names.append(f"治理與社會段（Step 3）：{govsoci_file.name}")
                else:
                    st.warning("⚠️ 未找到治理與社會段報告（Step 3）")
                
                if not files_to_merge:
                    st.error("❌ 沒有找到任何報告文件，請先完成 Step 1、2、3")
                    st.stop()
                
                # 顯示要合併的文件
                st.info("📋 **合併順序：**")
                for i, name in enumerate(file_names, 1):
                    st.write(f"{i}. {name}")
                
                # 生成輸出文件名
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_filename = f"ESG完整報告_{timestamp}.pptx"
                output_path = ESG_OUTPUT_ROOT / output_filename
                
                # 執行合併
                try:
                    total_slides = merge_pptx_files(files_to_merge, output_path)
                    
                    if output_path.exists():
                        st.success(f"✅ **彙整完成！**")
                        st.info(f"📁 **完整路徑：** `{output_path}`")
                        st.info(f"📊 **總頁數：** {total_slides} 頁")
                        st.session_state.step4_done = True
                        
                        # 下載按鈕
                        with open(output_path, "rb") as f:
                            file_data = f.read()
                        
                        st.download_button(
                            label="📥 下載完整 ESG 報告",
                            data=file_data,
                            file_name=output_filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            type="primary",
                            use_container_width=True
                        )
                        st.balloons()
                        
                        # 提示：如果文件有問題，可以手動合併
                        st.warning("""
                        ⚠️ **重要提示：** 
                        如果合併後的文件在 PowerPoint 中提示「需要修復」，這是正常的。
                        PowerPoint 會自動修復格式問題，點擊「修復」即可。
                        
                        **如果修復後仍有問題，建議使用手動合併：**
                        1. 打開 PowerPoint
                        2. 打開第一個報告（公司段）
                        3. 依次插入其他報告的投影片：
                           - 插入 → 投影片 → 重用投影片
                           - 選擇其他報告文件
                           - 選擇要插入的投影片
                        """)
                    else:
                        st.error(f"❌ 檔案儲存失敗！路徑：{output_path}")
                except Exception as merge_error:
                    st.error(f"❌ 自動合併失敗：{merge_error}")
                    st.warning("""
                    ⚠️ **自動合併遇到問題，建議手動合併：**
                    
                    **手動合併步驟：**
                    1. 打開 PowerPoint
                    2. 打開第一個報告：`{company_file.name if company_file else '公司段報告'}`
                    3. 依次插入其他報告：
                       - 插入 → 投影片 → 重用投影片
                       - 選擇 `{environment_file.name if environment_file else '環境段報告'}`
                       - 選擇 `{govsoci_file.name if govsoci_file else '治理與社會段報告'}`
                    4. 按順序插入所有投影片
                    5. 另存為完整報告
                    """.format(
                        company_file=company_file,
                        environment_file=environment_file,
                        govsoci_file=govsoci_file
                    ))
                    
                    # 提供個別文件的下載連結
                    st.subheader("📥 下載個別報告文件")
                    for file_path in files_to_merge:
                        if file_path.exists():
                            with open(file_path, "rb") as f:
                                file_data = f.read()
                            st.download_button(
                                label=f"📄 {file_path.name}",
                                data=file_data,
                                file_name=file_path.name,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key=f"download_{file_path.name}"
                            )
                    
            except Exception as e:
                st.error(f"❌ 合併失敗：{e}")
                st.exception(e)
else:
    st.warning("⚠️ 請先完成所有前置步驟")

