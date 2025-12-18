# audit_board.py
# -*- coding: utf-8 -*-
"""
Financial Audit Quarterly Board Generator (European Style)
----------------------------------------------------------
• 產出 A4 橫式 PowerPoint，2×2 區塊 (Qtr 1–4)
• 風格：歐洲印刷設計（霧面背板 + 淡灰導線 + 精簡對稱）
• 字體：Arial，跨平台穩定
• 預設內含審計自然語言文字，可直接呼叫
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def european_palette():
    """歐洲印刷風配色"""
    return {
        "title": (0, 51, 102),
        "text": (50, 50, 50),
        "panel_bg": (248, 249, 251),
        "panel_line": (214, 228, 240),
        "card_bg": (242, 245, 248),
        "heading": (0, 51, 102),
    }


def default_quarters():
    """預設審計內容（可修改成 ESG、治理等主題）"""
    return {
        "Qtr 1": (
            "第一季 — 基礎與風險基準線",
            "審計年度從確認期初餘額、驗證前一年度審計發現，以及調整年度審計計畫開始。"
            "內部稽核驗證會計、財務和薪資系統的控制作業，以確認職責分離和資料完整性。"
            "管理層的風險登記冊和重大性門檻會更新以對齊新的業務暴露。"
            "第一季建立驗證基準線，用於半年保證。"
        ),
        "Qtr 2": (
            "第二季 — 交易測試與控制執行",
            "審計抽樣擴展至採購、費用和日記帳分錄的交易。"
            "銀行和庫存對帳會獨立驗證，同時授權鏈會追溯至支持文件。"
            "例外情況會被記錄，矯正行動透過結案追蹤進行監控。"
            "第二季測試支持中期（半年）財務檢討，並確認年中控制可靠性。"
        ),
        "Qtr 3": (
            "第三季 — 流程驗證與治理檢討",
            "稽核人員對採購到付款 (P2P)、訂單到現金 (O2C) 和記錄到報告 (R2R) 循環進行流程層級評估，"
            "聚焦交易追溯性和會計準確性。"
            "公司間餘額、遞延稅項準備和 ESG 相關報告控制會測試一致性。"
            "年中報告整合發現、管理層回應和矯正行動進度。"
            "第三季確保在年度結帳階段前的營運控制連續性。"
        ),
        "Qtr 4": (
            "第四季 — 年終測試與最終保證",
            "審計重點轉移至調整的最終驗證、截止測試和財務結帳準確性。"
            "內部和外部稽核人員協調證據請求，以減少冗餘並強化文件。"
            "最終審計摘要概述控制有效性、未解決問題和下一年的風險優先順序。"
            "第四季為年度報告和治理簽核提供完整保證。"
        ),
    }


def render_audit_board(
    quarters=None,
    output_path="Financial_Audit_Qtr_Board.pptx",
    title="財務稽核實施計畫 — 精簡版",
    font_name="Arial",
    layout="A4L"
):
    """生成 PowerPoint 審計四季板"""

    # 預設資料與顏色
    quarters = quarters or default_quarters()
    colors = european_palette()

    prs = Presentation()
    if layout.upper() == "A4L":
        prs.slide_width = Inches(11.69)
        prs.slide_height = Inches(8.27)

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name = font_name
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(*colors["title"])

    # 背板
    panel = slide.shapes.add_shape(1, Inches(0.3), Inches(1.0), Inches(9.3), Inches(4.8))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(*colors["panel_bg"])
    panel.line.color.rgb = RGBColor(*colors["panel_line"])
    panel.line.width = Pt(0.75)

    # 位置
    card_positions = [
        (Inches(0.4), Inches(1.1)),
        (Inches(5.1), Inches(1.1)),
        (Inches(0.4), Inches(3.2)),
        (Inches(5.1), Inches(3.2)),
    ]

    # 生成四區塊
    for i, key in enumerate(["Qtr 1", "Qtr 2", "Qtr 3", "Qtr 4"]):
        if key not in quarters:
            continue
        title_text, body_text = quarters[key]
        left, top = card_positions[i]
        width, height = Inches(4.4), Inches(2.0)

        # 區塊底色
        card = slide.shapes.add_shape(1, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*colors["card_bg"])
        card.line.color.rgb = RGBColor(*colors["panel_line"])

        # 文字區
        textbox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.15), width - Inches(0.5), height - Inches(0.3))
        tf = textbox.text_frame
        tf.word_wrap = True

        # 標題
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.name = font_name
        p1.font.bold = True
        p1.font.size = Pt(11)
        p1.font.color.rgb = RGBColor(*colors["heading"])

        # 內容
        p2 = tf.add_paragraph()
        p2.text = body_text
        p2.font.name = font_name
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(*colors["text"])

    prs.save(output_path)
    return output_path


# ----------------------------------------------------------------------
# Component wrapper for embedding into existing slides
# ----------------------------------------------------------------------
class AuditBoardComponent:
    """
    Draw the quarterly audit board (4 cards) directly on a provided slide.
    """

    def __init__(
        self,
        prs,
        title="財務稽核實施計畫 — 精簡版",
        font_name="Arial",
        layout="A4L",
    ):
        self.prs = prs
        self.title = title
        self.font_name = font_name
        self.layout = layout
        self.colors = european_palette()

    def add_to_slide(
        self,
        slide,
        quarters=None,
        title=None,
        panel_left_in=0.3,
        panel_top_in=1.0,
        card_spacing_in=0.2,
    ):
        quarters = quarters or default_quarters()
        colors = self.colors
        title = title if title is not None else self.title

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.name = self.font_name
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(*colors["title"])

        # Panel
        panel_width = Inches(9.3)
        panel_height = Inches(4.8)
        panel = slide.shapes.add_shape(
            1,
            Inches(panel_left_in),
            Inches(panel_top_in),
            panel_width,
            panel_height,
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(*colors["panel_bg"])
        panel.line.color.rgb = RGBColor(*colors["panel_line"])
        panel.line.width = Pt(0.75)

        # Card positions (2x2 grid)
        card_width = Inches(4.4)
        card_height = Inches(2.0)
        horizontal_gap = Inches(card_spacing_in)
        vertical_gap = Inches(card_spacing_in + 1.1)

        base_left = Inches(panel_left_in + 0.1)
        base_top = Inches(panel_top_in + 0.1)

        positions = [
            (base_left, base_top),
            (base_left + card_width + horizontal_gap, base_top),
            (base_left, base_top + card_height + Inches(1.1)),
            (base_left + card_width + horizontal_gap, base_top + card_height + Inches(1.1)),
        ]

        order = ["Qtr 1", "Qtr 2", "Qtr 3", "Qtr 4"]

        for idx, key in enumerate(order):
            if key not in quarters:
                continue
            card_title, card_body = quarters[key]
            left, top = positions[idx]

            card = slide.shapes.add_shape(1, left, top, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*colors["card_bg"])
            card.line.color.rgb = RGBColor(*colors["panel_line"])

            textbox = slide.shapes.add_textbox(
                left + Inches(0.25),
                top + Inches(0.15),
                card_width - Inches(0.5),
                card_height - Inches(0.3),
            )
            tf = textbox.text_frame
            tf.word_wrap = True

            p1 = tf.add_paragraph()
            p1.text = card_title
            p1.font.name = self.font_name
            p1.font.bold = True
            p1.font.size = Pt(11)
            p1.font.color.rgb = RGBColor(*colors["heading"])

            p2 = tf.add_paragraph()
            p2.text = card_body
            p2.font.name = self.font_name
            p2.font.size = Pt(9)
            p2.font.color.rgb = RGBColor(*colors["text"])

        return slide


if __name__ == "__main__":
    path = render_audit_board()
    print(f"[OK] Generated file saved at: {path}")
