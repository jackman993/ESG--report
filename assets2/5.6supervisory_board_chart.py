
"""
Supervisory Board Chart Component (Modular Version)
--------------------------------------------------
用途：
    作為組織圖元件，可靈活嵌入任意 PowerPoint slide。
    不再輸出整頁，而是以座標方式畫出整個組織架構區塊。

可整合至：
    - 一鍵輸出模型
    - ESG/治理報告生成引擎
    - 自動簡報模組

使用：
    from supervisory_board_component import SupervisoryBoardChart
    chart = SupervisoryBoardChart(prs)
    chart.add_to_slide(slide, x=1.0, y=1.0)
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

class SupervisoryBoardChart:
    def __init__(self, prs, scale=0.8, font_scale=0.75):
        self.prs = prs
        self.scale = scale
        self.font_scale = font_scale

        self.COLORS = [
            RGBColor(25, 55, 165),
            RGBColor(0, 153, 153),
            RGBColor(0, 128, 255),
            RGBColor(128, 0, 128),
            RGBColor(51, 153, 51),
        ]

        self.DATA = [
            ("稽核委員會",
             "協助監督財務報告和風險管理控制的完整性與品質",
             "4 位成員"),
            ("ESG 委員會",
             "監督永續策略和績效，以創造長期價值",
             "3 位成員"),
            ("薪酬委員會",
             "制定和實施公平透明的薪酬政策",
             "4 位成員"),
            ("選任與提名委員會",
             "準備關鍵職位的標準和任命程序",
             "4 位成員"),
            ("技術委員會",
             "建議和監督數位轉型和技術策略",
             "4 位成員"),
        ]

    def add_to_slide(self, slide, x=1.0, y=1.0):
        """將組織圖畫到指定 slide 座標位置。"""
        scale = self.scale
        fscale = self.font_scale

        # Supervisory Board Title
        title_box = slide.shapes.add_textbox(
            Inches(x + 2.5 * scale), Inches(y + 0.1 * scale),
            Inches(3 * scale), Inches(0.5 * scale)
        )
        tf = title_box.text_frame
        tf.text = "稽核委員會"
        tf.paragraphs[0].font.size = Pt(28 * fscale)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Horizontal line
        line_y = y + 0.9 * scale
        line = slide.shapes.add_shape(1, Inches(x), Inches(line_y),
                                      Inches(9 * scale), Inches(0.03))
        fill = line.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Box Layout
        box_width = Inches(1.6 * scale)
        box_height = Inches(2.5 * scale)
        start_x = Inches(x)
        y_offset = Inches(y + 1.3 * scale)

        for i, (title, desc, members) in enumerate(self.DATA):
            color = self.COLORS[i]
            box_x = start_x + Inches(i * 1.9 * scale)

            # Committee Box
            shape = slide.shapes.add_shape(1, box_x, y_offset, box_width, box_height)
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = color

            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.07)
            tf.margin_bottom = Inches(0.07)

            p1 = tf.add_paragraph()
            p1.text = title
            p1.font.bold = True
            p1.font.size = Pt(15 * fscale)
            p1.font.color.rgb = RGBColor(255, 255, 255)

            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(11 * fscale)
            p2.font.color.rgb = RGBColor(255, 255, 255)
            p2.space_after = Pt(4)

            p3 = tf.add_paragraph()
            p3.text = members
            p3.font.bold = True
            p3.font.size = Pt(10 * fscale)
            p3.font.color.rgb = RGBColor(255, 255, 255)
            p3.alignment = PP_ALIGN.CENTER

            # Connector line
            center_x = box_x + box_width / 2
            vert_top = Inches(line_y)
            vert_height = y_offset - vert_top
            vert_line = slide.shapes.add_shape(1, center_x - Inches(0.01),
                                               vert_top, Inches(0.02), vert_height)
            vert_line.fill.solid()
            vert_line.fill.fore_color.rgb = RGBColor(0, 0, 0)

        return slide
