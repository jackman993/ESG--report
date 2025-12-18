# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def render_esg_compliance_table(
    rows,
    output_path="SME_ESG_Compliance_European_Print_Board.pptx",
    title="中小企業 ESG 遵循核心指標 — 歐洲印刷版",
    layout="A4L",  # "A4L" (A4 Landscape) or "16x9"
    font_name="Arial",
    colors=None,
    add_left_guide=True,
):
    """
    Render a European-print-style dense two-column ESG compliance table into a PPTX.

    Args:
        rows (list[tuple[str, str]] or list[dict]): data rows. Each row = (Category, Narrative)
            - If dicts, use keys {"category": ..., "narrative": ...}
        output_path (str): path to save the .pptx
        title (str): slide title
        layout (str): "A4L" for A4 landscape, or "16x9" for widescreen
        font_name (str): safe cross-platform font (e.g., "Arial")
        colors (dict): override style colors:
            {
              "page_bg": (248, 249, 251),
              "panel_border": (220, 224, 228),
              "header_bg": (234, 236, 238),
              "header_text": (46, 46, 46),
              "text": (59, 59, 59),
              "row_even": (242, 245, 248),
              "row_odd": (248, 249, 251),
              "guide": (214, 228, 240),
              "title": (46, 46, 46),
            }
        add_left_guide (bool): add subtle left alignment guide line (euro editorial feel)

    Returns:
        str: saved file path
    """
    # ---- default color palette (European print style, muted) ----
    palette = {
        "page_bg": (248, 249, 251),
        "panel_border": (220, 224, 228),
        "header_bg": (234, 236, 238),
        "header_text": (46, 46, 46),
        "text": (59, 59, 59),
        "row_even": (242, 245, 248),
        "row_odd": (248, 249, 251),
        "guide": (214, 228, 240),
        "title": (46, 46, 46),
    }
    if colors:
        palette.update(colors)

    prs = Presentation()

    # ---- page size ----
    if layout.upper() == "A4L":
        prs.slide_width = Inches(11.69)  # A4 landscape width
        prs.slide_height = Inches(8.27)  # A4 landscape height
    else:
        # default 16:9
        pass

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # ---- title ----
    title_shape = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), prs.slide_width - Inches(1.2), Inches(0.8))
    tf = title_shape.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = font_name
    p.font.color.rgb = RGBColor(*palette["title"])

    # ---- background panel (soft matte board) ----
    board_left = Inches(0.4)
    board_top = Inches(1.05)
    board_width = prs.slide_width - Inches(0.8)
    board_height = prs.slide_height - Inches(1.5)

    rect = slide.shapes.add_shape(
        autoshape_type_id=1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        left=board_left, top=board_top, width=board_width, height=board_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*palette["page_bg"])
    rect.line.color.rgb = RGBColor(*palette["panel_border"])
    rect.line.width = Pt(0.75)

    # ---- table ----
    headers = ["類別", "關鍵遵循敘述"]
    n_cols = 2
    n_rows = len(rows) + 1

    # inner margins for table
    table_left = board_left + Inches(0.2)
    table_top = board_top + Inches(0.1)
    table_width = board_width - Inches(0.4)
    table_height = board_height - Inches(0.2)

    table = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height).table

    # column widths (European proportion)
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = table_width - table.columns[0].width

    # header styling
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.size = Pt(12)
            para.font.name = font_name
            para.font.color.rgb = RGBColor(*palette["header_text"])
            para.alignment = PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*palette["header_bg"])

    # content rows
    def _get_text_pair(item):
        if isinstance(item, dict):
            return item.get("category", ""), item.get("narrative", "")
        return item[0], item[1]

    for r_idx, item in enumerate(rows, start=1):
        cat, narrative = _get_text_pair(item)

        # category cell
        c0 = table.cell(r_idx, 0)
        c0.text = cat
        for para in c0.text_frame.paragraphs:
            para.font.size = Pt(11.5)
            para.font.bold = True
            para.font.name = font_name
            para.font.color.rgb = RGBColor(*palette["text"])
            para.alignment = PP_ALIGN.LEFT

        # narrative cell (dense, long text)
        c1 = table.cell(r_idx, 1)
        c1.text = narrative
        for para in c1.text_frame.paragraphs:
            para.font.size = Pt(11)       # dense but readable
            para.font.name = font_name
            para.font.color.rgb = RGBColor(*palette["text"])
            para.alignment = PP_ALIGN.LEFT
            # fine spacing (pptx uses exact values; this is reasonably cross-platform)
            # para.line_spacing = 1.3  # optional; uncomment if you want a bit more breathing

        # alternating row fill (subtle)
        fill_rgb = palette["row_even"] if (r_idx % 2 == 0) else palette["row_odd"]
        c0.fill.solid(); c0.fill.fore_color.rgb = RGBColor(*fill_rgb)
        c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(*fill_rgb)

    # left alignment guide (euro editorial)
    if add_left_guide:
        guide = slide.shapes.add_shape(
            autoshape_type_id=1,
            left=board_left - Inches(0.1),
            top=table_top,
            width=Inches(0.02),
            height=table_height
        )
        guide.fill.solid()
        guide.fill.fore_color.rgb = RGBColor(*palette["guide"])
        guide.line.color.rgb = RGBColor(*palette["guide"])

    prs.save(output_path)
    return output_path


# ----------------------------------------------------------------------
# Component wrapper for integration with PPT engines (e.g. slide modules)
# ----------------------------------------------------------------------

DEFAULT_COMPONENT_ROWS = [
    {
        "category": "法規概覽",
        "narrative": (
            "建立並維護所有適用法規的最新登記冊（勞動、稅務、資料保護、環境）。"
            "為每個項目指派負責人和審查節奏，確保 100% 涵蓋率，並記錄不合規案例及矯正行動。"
            "自動化法規修正警示；進行年度法律審查並存檔證據。"
        ),
    },
    {
        "category": "財務與稅務誠信",
        "narrative": (
            "維持具完整追溯性的內部控制。職責分離、季度對帳，並由外部會計師驗證稅務申報。"
            "保留財務記錄 7 年，並標記證據以備稽核。"
        ),
    },
    {
        "category": "資料保護與資訊安全",
        "narrative": (
            "應用最小權限存取、多因素認證 (MFA) 和傳輸加密。遵循 PDPA/GDPR，維護處理活動記錄 (RoPA)，"
            "對高風險處理執行資料保護影響評估 (DPIA)，並進行年度滲透測試。保留違規日誌並執行零容忍洩漏政策。"
        ),
    },
    {
        "category": "供應鏈與倫理",
        "narrative": (
            "對關鍵供應商執行 ESG 盡職調查，將永續條款整合至合約中，並每年稽核高風險供應商。"
            "在永續溝通中揭露主要供應商績效和矯正行動。"
        ),
    },
    {
        "category": "吹哨者與治理",
        "narrative": (
            "提供匿名管道並保證不報復。30 天內結案並向董事會總結結果。"
            "進行年度治理有效性檢討並揭露政策更新。"
        ),
    },
]


class EuropeanComplianceTable:
    """
    Reusable component: draw the ESG compliance table directly on a given slide.
    """

    def __init__(
        self,
        prs,
        title="中小企業 ESG 遵循核心指標 — 歐洲印刷版",
        font_name="Arial",
        colors=None,
        add_left_guide=True,
    ):
        self.prs = prs
        self.title = title
        self.font_name = font_name
        self.colors = colors or {}
        self.add_left_guide = add_left_guide

        # reuse default palette from render function
        self.palette = {
            "page_bg": (248, 249, 251),
            "panel_border": (220, 224, 228),
            "header_bg": (234, 236, 238),
            "header_text": (46, 46, 46),
            "text": (59, 59, 59),
            "row_even": (242, 245, 248),
            "row_odd": (248, 249, 251),
            "guide": (214, 228, 240),
            "title": (46, 46, 46),
        }
        self.palette.update(self.colors)

    def add_to_slide(
        self,
        slide,
        rows=None,
        title=None,
        show_title=True,
        board_left_in=0.4,
        board_top_in=1.05,
        board_width_in=None,
        board_height_in=None,
        table_font_size=11,
    ):
        """
        Draw the table on the provided slide.

        Args:
            slide: pptx slide object
            rows: data rows (list of dict/tuple). Defaults to DEFAULT_COMPONENT_ROWS
            title: override title text
            show_title: whether to draw the title block
            board_*_in: positioning in inches (float)
        """
        rows = rows or DEFAULT_COMPONENT_ROWS
        rows = [
            r for r in rows
            if isinstance(r, dict) and r.get("category") != "吹哨者與治理"
        ]
        title = title if title is not None else self.title

        prs = self.prs
        board_left = Inches(board_left_in)
        board_top = Inches(board_top_in)
        board_width = (
            prs.slide_width - Inches(0.8)
            if board_width_in is None
            else Inches(board_width_in)
        )
        board_height = (
            prs.slide_height - Inches(1.5)
            if board_height_in is None
            else Inches(board_height_in)
        )

        if show_title and title:
            title_shape = slide.shapes.add_textbox(
                Inches(0.6),
                Inches(0.35),
                prs.slide_width - Inches(1.2),
                Inches(0.8),
            )
            tf = title_shape.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.name = self.font_name
            p.font.color.rgb = RGBColor(*self.palette["title"])

        rect = slide.shapes.add_shape(
            autoshape_type_id=1,
            left=board_left,
            top=board_top,
            width=board_width,
            height=board_height,
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(*self.palette["page_bg"])
        rect.line.color.rgb = RGBColor(*self.palette["panel_border"])
        rect.line.width = Pt(0.75)

        headers = ["類別", "關鍵遵循敘述"]
        n_cols = 2
        n_rows = len(rows) + 1

        table_left = board_left + Inches(0.2)
        table_top = board_top + Inches(0.1)
        table_width = board_width - Inches(0.4)
        table_height = board_height - Inches(0.2)

        table = slide.shapes.add_table(
            n_rows, n_cols, table_left, table_top, table_width, table_height
        ).table
 
        table.columns[0].width = Inches(1.4)
        table.columns[1].width = table_width - table.columns[0].width

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            for para in cell.text_frame.paragraphs:
                para.font.bold = True
                para.font.size = Pt(12)
                para.font.name = self.font_name
                para.font.color.rgb = RGBColor(*self.palette["header_text"])
                para.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*self.palette["header_bg"])

        for r_idx, item in enumerate(rows, start=1):
            if isinstance(item, dict):
                cat = item.get("category", "")
                narrative = item.get("narrative", "")
            else:
                cat, narrative = item[0], item[1]

            c0 = table.cell(r_idx, 0)
            c0.text = cat
            for para in c0.text_frame.paragraphs:
                para.font.size = Pt(11.5)
                para.font.bold = True
                para.font.name = self.font_name
                para.font.color.rgb = RGBColor(*self.palette["text"])
                para.alignment = PP_ALIGN.LEFT

            c1 = table.cell(r_idx, 1)
            c1.text = narrative
            for para in c1.text_frame.paragraphs:
                para.font.size = Pt(table_font_size)
                para.font.name = self.font_name
                para.font.color.rgb = RGBColor(*self.palette["text"])
                para.alignment = PP_ALIGN.LEFT

            fill_rgb = (
                self.palette["row_even"] if (r_idx % 2 == 0) else self.palette["row_odd"]
            )
            c0.fill.solid()
            c0.fill.fore_color.rgb = RGBColor(*fill_rgb)
            c1.fill.solid()
            c1.fill.fore_color.rgb = RGBColor(*fill_rgb)

        if self.add_left_guide:
            guide = slide.shapes.add_shape(
                autoshape_type_id=1,
                left=board_left - Inches(0.1),
                top=table_top,
                width=Inches(0.02),
                height=table_height,
            )
            guide.fill.solid()
            guide.fill.fore_color.rgb = RGBColor(*self.palette["guide"])
            guide.line.color.rgb = RGBColor(*self.palette["guide"])

        return slide


# ------------------------- Example usage -------------------------
if __name__ == "__main__":
    example_rows = [
        {"category": "法規概覽",
         "narrative": ("建立並維護所有適用法規的最新登記冊（勞動、稅務、資料保護、環境）。"
                       "為每個項目指派負責人和審查節奏，確保 100% 涵蓋率，並記錄不合規案例及矯正行動。"
                       "自動化法規修正警示；進行年度法律審查並存檔證據。")},
        {"category": "財務與稅務誠信",
         "narrative": ("維持具完整追溯性的內部控制。職責分離、季度對帳，並由外部會計師驗證稅務申報。"
                       "保留財務記錄 7 年，並標記證據以備稽核。")},
        {"category": "資料保護與資訊安全",
         "narrative": ("應用最小權限存取、多因素認證 (MFA) 和傳輸加密。遵循 PDPA/GDPR，維護處理活動記錄 (RoPA)，"
                       "對高風險處理執行資料保護影響評估 (DPIA)，並進行年度滲透測試。保留違規日誌並執行零容忍洩漏政策。")},
        {"category": "供應鏈與倫理",
         "narrative": ("對關鍵供應商執行 ESG 盡職調查，將永續條款整合至合約中，並每年稽核高風險供應商。"
                       "在永續溝通中揭露主要供應商績效和矯正行動。")},
        {"category": "吹哨者與治理",
         "narrative": ("提供匿名管道並保證不報復。30 天內結案並向董事會總結結果。"
                       "進行年度治理有效性檢討並揭露政策更新。")},
    ]

    path = render_esg_compliance_table(
        rows=example_rows,
        output_path="SME_ESG_Compliance_European_Print_Board.pptx",
        layout="A4L",
        font_name="Arial",
    )
    print(f"Saved: {path}")
