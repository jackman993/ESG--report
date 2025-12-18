"""
Board Remuneration table component.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


class BoardIncentiveFramework:
    def __init__(self, prs, scale=1.0):
        self.prs = prs
        self.scale = scale
        self.blocks = [
            {
                "title": "短期激勵 (STI)",
                "rationale": "鼓勵年度財務、營運和創新目標的達成，同時維持競爭力。",
                "policy": [
                    "年度績效現金獎金",
                    "指標包括財務結果、技術里程碑、客戶滿意度",
                    "由薪酬委員會每年審查"
                ],
            },
            {
                "title": "長期激勵 (LTI)",
                "rationale": "將高階主管與長期股東價值和永續目標對齊。",
                "policy": [
                    "多年歸屬的股權獎勵",
                    "目標包括總股東報酬率 (TSR)、ESG 指標、技術領導地位",
                    "由監事會核准以強化長期價值"
                ],
            },
            {
                "title": "退休金與福利",
                "rationale": "提供符合市場標準的補充福利，支持留任和福祉。",
                "policy": [
                    "確定提撥退休金和補充計畫",
                    "報銷、津貼和保險保障"
                ],
            },
            {
                "title": "持股指引",
                "rationale": "透過最低持股要求，強化管理層與股東的對齊。",
                "policy": [
                    "執行長：3 倍基本薪資；其他管理委員會成員：2 倍基本薪資",
                    "五年合規窗口期",
                    "臨時偏差需經監事會核准"
                ],
            },
        ]

    def add_to_slide(self, slide, x=0.7, y=1.2):
        rows = len(self.blocks) + 1
        cols = 3
        table_width = Inches(11.5 * self.scale)
        table_height = Inches(4.0 * self.scale)
        left = Inches(x)
        top = Inches(y)

        table = slide.shapes.add_table(rows, cols, left, top, table_width, table_height).table
        col_widths = [Inches(3.0 * self.scale), Inches(4.0 * self.scale), Inches(4.5 * self.scale)]
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = width

        headers = ["組成項目", "策略連結 / 理由", "政策 / 機制摘要"]
        for idx, header in enumerate(headers):
            cell = table.cell(0, idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(96, 96, 96)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.name = "Calibri"

        for row_idx, block in enumerate(self.blocks, start=1):
            # Component title
            cell = table.cell(row_idx, 0)
            cell.text = block["title"]
            self._format_cell(cell, bold=True)

            # Rationale
            cell = table.cell(row_idx, 1)
            cell.text = block["rationale"]
            self._format_cell(cell)

            # Policy list
            cell = table.cell(row_idx, 2)
            policy_text = "\n".join(f"• {item}" for item in block["policy"])
            cell.text = policy_text
            self._format_cell(cell)

        return slide

    def _format_cell(self, cell, bold=False):
        cell.text_frame.word_wrap = True
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.name = "Calibri"
            paragraph.font.bold = bold
            paragraph.font.color.rgb = RGBColor(34, 34, 34)
            paragraph.space_after = Pt(4)
