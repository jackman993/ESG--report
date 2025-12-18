"""
檢查母版選擇邏輯
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 模板路徑
SEED_TEMPLATE_PATH = r"C:\Users\User\Desktop\handdrawppt.pptx"

print("=" * 60)
print("檢查母版選擇邏輯")
print("=" * 60)
print()

if not Path(SEED_TEMPLATE_PATH).exists():
    print(f"❌ 模板不存在: {SEED_TEMPLATE_PATH}")
    sys.exit(1)

# 載入模板
prs = Presentation(SEED_TEMPLATE_PATH)
layouts = prs.slide_layouts

print(f"📄 模板路徑: {SEED_TEMPLATE_PATH}")
print(f"📊 模板現有頁面數: {len(prs.slides)}")
print(f"📋 可用母版數: {len(layouts)}")
print()

# 列出所有母版
print("所有可用的母版:")
for idx, layout in enumerate(layouts):
    name = layout.name if hasattr(layout, 'name') else f'Layout {idx}'
    print(f"  [{idx}] {name}")
print()

# 檢查模板中現有頁面的母版
if len(prs.slides) > 0:
    print("模板中現有頁面的母版:")
    for idx, slide in enumerate(prs.slides):
        layout = slide.slide_layout
        layout_name = layout.name if hasattr(layout, 'name') else 'Unknown'
        
        # 計算 placeholder 數量
        placeholder_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                placeholder_count += 1
        
        print(f"  第 {idx + 1} 頁: {layout_name} (有 {placeholder_count} 個 placeholder)")
    print()

# 選擇邏輯（模擬程式中的選擇）
print("程式選擇邏輯:")
print("-" * 60)

# 使用最後一個 layout（程式中的邏輯）
if len(layouts) > 0:
    selected_layout = layouts[len(layouts) - 1]
    selected_name = selected_layout.name if hasattr(selected_layout, 'name') else f'Layout {len(layouts)-1}'
    print(f"✅ 選擇的母版: {selected_name} (最後一個)")
    print(f"   索引: {len(layouts) - 1}")
else:
    print("❌ 沒有可用的母版")

print()
print("=" * 60)
print("檢查完成")
print("=" * 60)

