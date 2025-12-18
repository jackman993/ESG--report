"""Company PPT engine supporting layouts A, B, and C."""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Any, Optional, Tuple
import re
import importlib.util
import os
from pathlib import Path
import sys

from config_pptx_company import PPT_CONFIG, SLIDE_CONFIGS, SEED_TEMPLATE_PATH, OUTPUT_PATH
from content_pptx_company import PPTContentEngine

# 移除 auto_repair_pptx 調用（修正引擎沒有用，不需要調度）
# auto_repair_pptx = None

CM_TO_INCH = 1 / 2.54


def cm(value: float) -> float:
    return value * CM_TO_INCH


class PPTFullEngine:
    def __init__(self, content_engine: PPTContentEngine, company_name: str = None):
        if not os.path.exists(SEED_TEMPLATE_PATH):
            raise FileNotFoundError(f"Seed template missing: {SEED_TEMPLATE_PATH}")
        self.prs = Presentation(SEED_TEMPLATE_PATH)
        self.content_engine = content_engine
        self.slide_configs = SLIDE_CONFIGS
        self.output_path = Path(OUTPUT_PATH)
        self.output_path.mkdir(exist_ok=True)
        # 統一使用 Microsoft JhengHei 作為預設字型（含中英文），
        # 由 PowerPoint 負責英文 fallback 到 Calibri，避免使用怪字體。
        self.font_family = PPT_CONFIG.get("font_family", "Microsoft JhengHei")
        self.company_name = company_name  # 公司名稱，用於替換

        # 偵錯用開關：
        # - DISABLE_IMAGES=1  -> 不插入任何圖片（add_picture 不執行）
        # - DISABLE_COMPONENTS=1 -> 不渲染任何 component（包含 table / flowchart 等）
        self.disable_images = os.getenv("DISABLE_IMAGES", "0") == "1"
        self.disable_components = os.getenv("DISABLE_COMPONENTS", "0") == "1"
        # - DISABLE_FLOWS=1 -> 關閉流程圖類 component（例如 RiskFlowchartComponent）
        self.disable_flows = os.getenv("DISABLE_FLOWS", "0") == "1"

        # 預先建立所有投影片（重用模板前5頁作為前5張內容頁，其餘以 add_slide 補足到 17 張）
        self.total_slides = PPT_CONFIG.get("total_slides", 17)
        existing_slides = len(self.prs.slides)  # 模板原有頁數（預期為 5）
        needed_slides = self.total_slides

        self._slides_by_index = {}
        layouts = self.prs.slide_layouts
        # 優先使用第一個版面（在 minimal 測試中確認過是安全的）
        base_layout = layouts[0]

        # 先將現有模板頁面映射為前幾張內容頁（例如 1~5）
        for idx in range(1, min(existing_slides, needed_slides) + 1):
            self._slides_by_index[idx] = self.prs.slides[idx - 1]

        # 如有需要，再以 add_slide 補足到 total_slides
        for idx in range(existing_slides + 1, needed_slides + 1):
            new_slide = self.prs.slides.add_slide(base_layout)
            self._slides_by_index[idx] = new_slide

        print(
            f"[INFO] 預先建立 {needed_slides} 張投影片完成（重用模板原有 {existing_slides} 頁，最終總頁數 {len(self.prs.slides)}）"
        )

    def generate(self, company_name: str = None):
        """生成 PPT，可選的公司名稱用於替換 CEO message 中的 'our company'"""
        if company_name:
            self.company_name = company_name
        for idx in range(1, self.total_slides + 1):
            self._generate_slide(idx)

        # 基礎檔名
        base_name = PPT_CONFIG.get("output_filename", "ESG_PPT_company.pptx")
        base = Path(base_name)

        # 為避免檔案被鎖定，每次都帶上時間戳
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output = self.output_path / f"{base.stem}_{timestamp}{base.suffix}"

        try:
            self.prs.save(str(output))
            print(f"[OK] PPT saved -> {output}")
        except PermissionError as e:
            # 若仍遇到鎖檔，再試一次用不同檔名
            print(f"[WARN] 儲存檔案時發生 PermissionError，嘗試使用新檔名: {e}")
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output = self.output_path / f"{base.stem}_{timestamp}_alt{base.suffix}"
            self.prs.save(str(output))
            print(f"[OK] PPT saved with alternate name -> {output}")

        # 已完全禁用 auto_repair_pptx 以避免 UI 卡頓
        # 移除 auto_repair_pptx 調用（修正引擎沒有用，不需要調度）
        # if auto_repair_pptx is not None:
        #     try:
        #         fixed_path = auto_repair_pptx(str(output))
        #         output = Path(fixed_path)
        #     except Exception as e:
        #         print(f"[WARN] 自動修復公司段 PPT 失敗，將使用未修復檔案：{e}")

        return str(output)

    def generate_subset(self, slide_indices: List[int], output_filename: Optional[str] = None):
        for idx in slide_indices:
            self._generate_slide(idx)
        filename = output_filename or PPT_CONFIG.get("output_filename", "ESG_PPT_company_experimental.pptx")
        output = self.output_path / filename
        self.prs.save(str(output))
        print(f"[OK] PPT subset saved -> {output}")
        return str(output)

    def _generate_slide(self, slide_index: int):
        cfg = self.slide_configs[slide_index]
        slide = self._ensure_slide(slide_index)
        print(f"[Slide {slide_index}] {cfg.get('title','Untitled')}")
        # 不再清空整張投影片的 XML 結構，只是在預先建立好的空白頁上填內容
        self._apply_title(slide, cfg)
        layout = cfg.get("layout")
        if layout == "A":
            self._layout_a(slide, cfg)
        elif layout == "B":
            self._layout_b(slide, cfg)
        elif layout == "C":
            self._layout_c(slide, cfg)
        else:
            raise ValueError(f"Unsupported layout: {layout}")
        self._add_watermark(slide, PPT_CONFIG.get("watermark_text"))

    def _ensure_slide(self, slide_index: int):
        """回傳預先建立好的第 slide_index 張投影片。"""
        slide = self._slides_by_index.get(slide_index)
        if slide is None:
            raise IndexError(f"Slide index out of range: {slide_index}")
        return slide

    def _apply_title(self, slide, cfg):
        title = cfg.get("title", "")
        if not title:
            return
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and getattr(shape.placeholder_format, "type", None) == 1:
                shape.text = title
                shape.top = Inches(cm(cfg.get("title_top_cm", 1.5)))
                shape.left = Inches(cm(cfg.get("title_left_cm", 2.5)))
                shape.width = Inches(cm(cfg.get("title_width_cm", 20.0)))
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = self.font_family
                    paragraph.font.size = Pt(cfg.get("title_font_pt", 16))
                    paragraph.font.bold = True
                return
        box = slide.shapes.add_textbox(
            Inches(cm(cfg.get("title_left_cm", 2.5))),
            Inches(cm(cfg.get("title_top_cm", 1.5))),
            Inches(cm(cfg.get("title_width_cm", 20.0))),
            Inches(1.0),
        )
        box.text_frame.text = title
        p = box.text_frame.paragraphs[0]
        p.font.name = self.font_family
        p.font.size = Pt(cfg.get("title_font_pt", 16))
        p.font.bold = True

    def _layout_a(self, slide, cfg):
        paragraphs = cfg.get("paragraphs", 2)
        if paragraphs:
            left_box_cfg = cfg.get("left_text_box", {})
            text_box = self._get_textbox(
                slide,
                left_cm=left_box_cfg.get("left_cm", 2.3),
                top_cm=left_box_cfg.get("top_cm", 3.0),
                width_cm=left_box_cfg.get("width_cm", 13.5),
                height_cm=left_box_cfg.get("height_cm", 15.0),
            )
            text = self._get_slide_text(cfg)
            self._fill_textbox(
                text_box,
                text,
                paragraphs,
                cfg.get("text_font_size", 11),
                spacing_before_pt=cfg.get("paragraph_spacing_pt"),
                spacing_after_pt=cfg.get("paragraph_spacing_pt"),
                font_color_rgb=cfg.get("text_color_rgb"),
                placeholders=cfg.get("left_placeholders"),
            )
            self._append_notes(text_box, cfg.get("left_notes"))
            if cfg.get("left_alert_text"):
                self._apply_alert(text_box, cfg["left_alert_text"])

        left_component_cfg = cfg.get("left_component")
        if left_component_cfg and not self.disable_components:
            self._render_component(slide, left_component_cfg)

        background_cfg = cfg.get("right_background")
        if background_cfg and not self.disable_images:
            path = background_cfg.get("image_path")
            if path and Path(path).exists():
                left_cm_val = background_cfg.get("left_cm", 18.0)
                top_cm = background_cfg.get("top_cm", 3.0)
                width_cm_val = background_cfg.get("width_cm")
                height_cm_val = background_cfg.get("height_cm")
                pic = slide.shapes.add_picture(
                    str(path),
                    Inches(cm(left_cm_val)),
                    Inches(cm(top_cm)),
                )
                if width_cm_val:
                    pic.width = Inches(cm(width_cm_val))
                if height_cm_val:
                    pic.height = Inches(cm(height_cm_val))
                if background_cfg.get("send_to_back"):
                    slide.shapes._spTree.remove(pic._element)
                    slide.shapes._spTree.insert(0, pic._element)

        image_path = cfg.get("image_path")
        component_cfg = cfg.get("right_component")
        if component_cfg:
            if not self.disable_components:
                self._render_component(slide, component_cfg)
        elif image_path and Path(image_path).exists() and not self.disable_images:
            # ⚠ 圖片插入策略（簡化版，對齊環境段邏輯）：
            # - 不再在插入後動態修改 width/height / left
            # - 直接在 add_picture 時給定寬高，讓 PowerPoint 自己處理長寬比
            width_cm = cfg.get("image_width_cm")
            height_cm = cfg.get("image_height_cm")
            top_cm = cfg.get("image_top_cm", 4.0)
            left_cm_val = cfg.get("image_left_cm", 19.0)

            picture_kwargs = {}
            if width_cm is not None and height_cm is not None:
                # 同時給寬高，邏輯跟 environment_pptx._add_image_full_path 類似
                picture_kwargs["width"] = Inches(cm(width_cm))
                picture_kwargs["height"] = Inches(cm(height_cm))
            elif width_cm is not None:
                picture_kwargs["width"] = Inches(cm(width_cm))
            elif height_cm is not None:
                picture_kwargs["height"] = Inches(cm(height_cm))
            else:
                picture_kwargs["width"] = Inches(cm(14.0))

            pic = slide.shapes.add_picture(
                str(image_path),
                Inches(cm(left_cm_val)),
                Inches(cm(top_cm)),
                **picture_kwargs,
            )

            # 不做額外的縮放 / 右對齊計算，先確認這種「環境段式」寫法是否還會觸發修復
            if cfg.get("image_alert_text"):
                self._apply_alert(pic, cfg["image_alert_text"])

        right_paragraphs = cfg.get("right_paragraphs", 0)
        if right_paragraphs:
            right_box_cfg = cfg.get("right_text_box", {})
            right_box = self._get_textbox(
                slide,
                left_cm=right_box_cfg.get("left_cm", 18.0),
                top_cm=right_box_cfg.get("top_cm", 3.0),
                width_cm=right_box_cfg.get("width_cm", 13.5),
                height_cm=right_box_cfg.get("height_cm", 15.0),
            )
            right_text = self._get_slide_text(cfg, prefix="right")
            self._fill_textbox(
                right_box,
                right_text,
                right_paragraphs,
                cfg.get("right_text_font_size", cfg.get("text_font_size", 11)),
                spacing_before_pt=cfg.get("right_paragraph_spacing_pt", cfg.get("paragraph_spacing_pt")),
                spacing_after_pt=cfg.get("right_paragraph_spacing_pt", cfg.get("paragraph_spacing_pt")),
                font_color_rgb=cfg.get("right_text_color_rgb") or cfg.get("text_color_rgb"),
                placeholders=cfg.get("right_placeholders"),
            )
            self._append_notes(right_box, cfg.get("right_notes"))
            if cfg.get("right_alert_text"):
                self._apply_alert(right_box, cfg["right_alert_text"])

    def _layout_b(self, slide, cfg):
        text_box = None
        paragraphs = cfg.get("paragraphs", 2)
        if paragraphs:
            text_box_cfg = cfg.get("text_box", {})
            text_box = self._get_textbox(
                slide,
                left_cm=text_box_cfg.get("left_cm", 17.2),
                top_cm=text_box_cfg.get("top_cm", 3.2),
                width_cm=text_box_cfg.get("width_cm", 15.0),
                height_cm=text_box_cfg.get("height_cm", 14.5),
            )
            text = self._get_slide_text(cfg)
            self._fill_textbox(
                text_box,
                text,
                paragraphs,
                cfg.get("text_font_size", 11),
                spacing_before_pt=cfg.get("paragraph_spacing_pt"),
                spacing_after_pt=cfg.get("paragraph_spacing_pt"),
                font_color_rgb=cfg.get("text_color_rgb"),
                placeholders=cfg.get("left_placeholders"),
            )
            self._append_notes(text_box, cfg.get("right_notes") or cfg.get("text_box_notes"))
            if cfg.get("left_alert_text"):
                self._apply_alert(text_box, cfg["left_alert_text"])

        media_cfg = cfg.get("left_media", {})
        mode = media_cfg.get("mode", "stack")
        if mode == "component":
            if not self.disable_components:
                self._render_component(slide, media_cfg)
            return

        paths = cfg.get("image_paths", [])
        if not paths:
            return
        area_left = media_cfg.get("area_left_cm", 3.0)
        area_top = media_cfg.get("area_top_cm", 3.5)
        area_height = media_cfg.get("area_height_cm", 13.0)
        overlay = media_cfg.get("overlay_box")
        if overlay:
            self._draw_overlay(slide, overlay)

        if mode == "stack":
            width_cm_val = media_cfg.get("width_cm", 8.0)
            gap_cm = media_cfg.get("gap_cm", 3.0)
            current_top = area_top
            area_right = area_left + width_cm_val
            pictures = []
            for path in paths:
                if not Path(path).exists() or self.disable_images:
                    continue
                pic = slide.shapes.add_picture(str(path), Inches(cm(area_left)), Inches(cm(current_top)), width=Inches(cm(width_cm_val)))
                pic.left = Inches(cm(area_right)) - pic.width
                pictures.append(pic)
                current_top += pic.height * CM_TO_INCH + gap_cm
            if pictures:
                used_height = sum(pic.height for pic in pictures) + Inches(cm(gap_cm)) * (len(pictures) - 1)
                available = Inches(cm(area_height))
                if used_height < available:
                    offset = (available - used_height) / 2
                    baseline = Inches(cm(area_top)) + offset
                    for idx, pic in enumerate(pictures):
                        pic.top = int(baseline)
                        baseline += pic.height
                        if idx < len(pictures) - 1:
                            baseline += Inches(cm(gap_cm))
        elif mode == "gallery":
            width_cm_val = media_cfg.get("width_cm", 14.0)
            widths_override = media_cfg.get("widths_cm", [])
            gap_cm = media_cfg.get("gap_cm", 1.0)
            caption_font = media_cfg.get("caption_font_size", 10)
            caption_height_cm = media_cfg.get("caption_height_cm", 1.0)
            caption_offset_cm = media_cfg.get("caption_offset_cm", 0.2)
            layout_variant = media_cfg.get("layout")
            captions = cfg.get("captions", [])
            gallery_items = []
            if layout_variant == "left_stack_right_single" and len(paths) >= 2 and not self.disable_images:
                left_paths = paths[:-1]
                right_path = paths[-1]
                left_items = []
                left_widths = []
                for idx, path in enumerate(left_paths):
                    if not Path(path).exists():
                        continue
                    width_cm_current = widths_override[idx] if idx < len(widths_override) else width_cm_val
                    pic = slide.shapes.add_picture(
                        str(path),
                        Inches(cm(area_left)),
                        Inches(cm(area_top)),
                        width=Inches(cm(width_cm_current)),
                    )
                    caption_box = None
                    if idx < len(captions) and captions[idx]:
                        caption_box = slide.shapes.add_textbox(
                            pic.left,
                            0,
                            pic.width,
                            Inches(cm(caption_height_cm)),
                        )
                        frame = caption_box.text_frame
                        frame.clear()
                        caption_para = frame.paragraphs[0]
                        caption_para.text = captions[idx]
                        caption_para.font.name = self.font_family
                        caption_para.font.size = Pt(caption_font)
                        caption_para.alignment = PP_ALIGN.CENTER
                        caption_para.space_before = Pt(0)
                        caption_para.space_after = Pt(0)
                    if caption_box:
                        caption_box.fill.background()
                    left_items.append((pic, caption_box))
                    left_widths.append(width_cm_current)

                if not Path(right_path).exists():
                    return
                right_width_cm = widths_override[len(paths) - 1] if len(widths_override) >= len(paths) else width_cm_val
                right_pic = slide.shapes.add_picture(
                    str(right_path),
                    Inches(cm(area_left)),
                    Inches(cm(area_top)),
                    width=Inches(cm(right_width_cm)),
                )
                right_caption = None
                if len(captions) >= len(paths) and captions[-1]:
                    right_caption = slide.shapes.add_textbox(
                        right_pic.left,
                        0,
                        right_pic.width,
                        Inches(cm(caption_height_cm)),
                    )
                    frame = right_caption.text_frame
                    frame.clear()
                    caption_para = frame.paragraphs[0]
                    caption_para.text = captions[-1]
                    caption_para.font.name = self.font_family
                    caption_para.font.size = Pt(caption_font)
                    caption_para.alignment = PP_ALIGN.CENTER
                    caption_para.space_before = Pt(0)
                    caption_para.space_after = Pt(0)

                left_gap = Inches(cm(gap_cm))
                caption_offset = Inches(cm(caption_offset_cm))
                available_left = Inches(cm(area_height)) if area_height else None
                left_total = 0
                for idx, (pic, caption_box) in enumerate(left_items):
                    left_total += pic.height
                    if caption_box:
                        left_total += caption_offset + caption_box.height
                    if idx < len(left_items) - 1:
                        left_total += left_gap
                if available_left is None:
                    available_left = left_total
                start_left = Inches(cm(area_top))
                if left_total < available_left:
                    start_left += (available_left - left_total) / 2

                cursor_left = start_left
                for idx, (pic, caption_box) in enumerate(left_items):
                    pic.left = Inches(cm(area_left))
                    pic.top = int(cursor_left)
                    cursor_left += pic.height
                    if caption_box:
                        caption_box.left = pic.left
                        caption_box.width = pic.width
                        caption_box.top = int(cursor_left + caption_offset)
                        cursor_left = caption_box.top + caption_box.height
                    if idx < len(left_items) - 1:
                        cursor_left += left_gap

                column_gap_cm = media_cfg.get("column_gap_cm", gap_cm)
                max_left_width = max(left_widths) if left_widths else width_cm_val
                right_left_cm = media_cfg.get("right_left_cm", area_left + max_left_width + column_gap_cm)
                available_right = Inches(cm(area_height)) if area_height else None
                right_total = right_pic.height
                if right_caption:
                    right_total += caption_offset + right_caption.height
                if available_right is None:
                    available_right = right_total
                start_right = Inches(cm(area_top))
                if right_total < available_right:
                    start_right += (available_right - right_total) / 2
                right_pic.left = Inches(cm(right_left_cm))
                right_pic.top = int(start_right)
                if right_caption:
                    right_caption.left = right_pic.left
                    right_caption.width = right_pic.width
                    right_caption.top = int(right_pic.top + right_pic.height + caption_offset)
                return

            for idx, path in enumerate(paths):
                if not Path(path).exists():
                    continue
                width_cm_current = widths_override[idx] if idx < len(widths_override) else width_cm_val
                pic = slide.shapes.add_picture(
                    str(path),
                    Inches(cm(area_left)),
                    Inches(cm(area_top)),
                    width=Inches(cm(width_cm_current)),
                )
                caption_box = None
                caption_height = 0
                if idx < len(captions) and captions[idx]:
                    caption_box = slide.shapes.add_textbox(
                        pic.left,
                        0,
                        pic.width,
                        Inches(cm(caption_height_cm)),
                    )
                    frame = caption_box.text_frame
                    frame.clear()
                    caption_para = frame.paragraphs[0]
                    caption_para.text = captions[idx]
                    caption_para.font.name = self.font_family
                    caption_para.font.size = Pt(caption_font)
                    caption_para.alignment = PP_ALIGN.CENTER
                    caption_para.space_before = Pt(0)
                    caption_para.space_after = Pt(0)
                    caption_height = caption_box.height
                    if caption_box:
                        caption_box.fill.background()
                gallery_items.append((pic, caption_box, pic.height + caption_height))

            if not gallery_items:
                return

            total_height = sum(item[2] for item in gallery_items) + Inches(cm(gap_cm)) * (len(gallery_items) - 1)
            available = Inches(cm(area_height)) if area_height else total_height
            start_top = Inches(cm(area_top))
            if total_height < available:
                start_top += (available - total_height) / 2

            cursor = start_top
            for idx, (pic, caption_box, _) in enumerate(gallery_items):
                pic.top = int(cursor)
                cursor += pic.height
                if caption_box:
                    caption_box.top = int(cursor)
                    cursor += caption_box.height
                if idx < len(gallery_items) - 1:
                    cursor += Inches(cm(gap_cm))
            if cfg.get("image_alert_text"):
                for pic, _, _ in gallery_items:
                    self._apply_alert(pic, cfg["image_alert_text"])
        elif mode == "single":
            width_cm_val = media_cfg.get("width_cm", 14.0)
            top_cm = media_cfg.get("area_top_cm", area_top)
            left_cm_val = media_cfg.get("area_left_cm", area_left)
            image_path = paths[0]
            if Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), Inches(cm(left_cm_val)), Inches(cm(top_cm)), width=Inches(cm(width_cm_val)))

    def _layout_c(self, slide, cfg):
        top_text = cfg.get("top_text") or self._get_slide_text(cfg)
        bottom_text = cfg.get("bottom_text")

        if top_text:
            top_box = self._get_textbox(
                slide,
                left_cm=cfg.get("top_box_left_cm", 2.5),
                top_cm=cfg.get("top_box_top_cm", 3.0),
                width_cm=cfg.get("top_box_width_cm", 22.0),
                height_cm=cfg.get("top_box_height_cm", 5.0),
            )
            self._fill_textbox(
                top_box,
                top_text,
                cfg.get("top_paragraphs", 2),
                cfg.get("text_font_size", 11),
                font_color_rgb=cfg.get("text_color_rgb"),
                placeholders=cfg.get("top_placeholders"),
            )
            self._append_notes(top_box, cfg.get("top_notes"))
            if cfg.get("top_alert_text"):
                self._apply_alert(top_box, cfg["top_alert_text"])

        media_cfg = cfg.get("center_media", {})
        mode = media_cfg.get("mode", "component")
        top_offset_cm = media_cfg.get("top_cm", 8.5)

        if mode == "component":
            self._render_component(
                slide,
                {
                    "file": media_cfg.get("file"),
                    "class": media_cfg.get("class"),
                    "method": media_cfg.get("method", "add_to_slide"),
                    "init_kwargs": media_cfg.get("init_kwargs", {}),
                    "method_kwargs": media_cfg.get("method_kwargs", {}),
                },
            )
        elif mode == "image":
            image_path = media_cfg.get("path")
            width_cm = media_cfg.get("width_cm", 14.0)
            left_cm_val = media_cfg.get("left_cm", 2.5)
            if image_path and Path(image_path).exists():
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(cm(left_cm_val)),
                    Inches(cm(top_offset_cm)),
                    width=Inches(cm(width_cm)),
                )

        if bottom_text:
            bottom_box = self._get_textbox(
                slide,
                left_cm=cfg.get("bottom_box_left_cm", 2.5),
                top_cm=media_cfg.get("bottom_text_top_cm", 15.5),
                width_cm=cfg.get("bottom_box_width_cm", 22.0),
                height_cm=cfg.get("bottom_box_height_cm", 4.8),
            )
            self._fill_textbox(
                bottom_box,
                bottom_text,
                cfg.get("bottom_paragraphs", 2),
                cfg.get("text_font_size", 11),
                font_color_rgb=cfg.get("text_color_rgb"),
                placeholders=cfg.get("bottom_placeholders"),
            )
            self._append_notes(bottom_box, cfg.get("bottom_notes"))
            if cfg.get("bottom_alert_text"):
                self._apply_alert(bottom_box, cfg["bottom_alert_text"])

    # Helpers (same as base engine)
    def _get_slide_text(self, cfg: Dict[str, Any], prefix: str = "main") -> str:
        if prefix in ("main", "left"):
            method_key = "content_method"
            fallback_key = "fallback_text"
        else:
            method_key = f"{prefix}_content_method"
            fallback_key = f"{prefix}_fallback_text"

        method_name = cfg.get(method_key)
        if method_name and hasattr(self.content_engine, method_name):
            text = getattr(self.content_engine, method_name)()
        else:
            text = cfg.get(fallback_key, "")
        
        # 替換佔位符
        # 如果有公司名稱，使用實際名稱；否則使用 "本公司"
        replacement_name = self.company_name if self.company_name else "本公司"
        
        # 替換 {COMPANY_NAME} 佔位符
        text = text.replace("{COMPANY_NAME}", replacement_name)
        
        # 如果是 CEO message，也替換 "our company" 為實際公司名稱
        if method_name == "generate_ceo_message":
            if self.company_name:
                text = text.replace("our company", self.company_name)
                text = text.replace("Our company", self.company_name)
            else:
                text = text.replace("our company", "本公司")
                text = text.replace("Our company", "本公司")
        
        return text

    def _get_textbox(self, slide, left_cm: float, top_cm: float, width_cm: float, height_cm: float):
        box = slide.shapes.add_textbox(Inches(cm(left_cm)), Inches(cm(top_cm)), Inches(cm(width_cm)), Inches(cm(height_cm)))
        box.text_frame.word_wrap = True
        return box

    def _fill_textbox(
        self,
        textbox,
        text: str,
        paragraphs: int,
        font_size: int,
        spacing_before_pt: Any = None,
        spacing_after_pt: Any = None,
        font_color_rgb: Optional[Tuple[int, int, int]] = None,
        placeholders: Optional[Dict[str, Dict[str, Any]]] = None,
    ):
        frame = textbox.text_frame
        frame.clear()
        if not text:
            frame.text = ""
            return
        chunks = self._split_into_paragraphs(text, max(paragraphs, 1))
        default_color = font_color_rgb or (58, 58, 58)
        placeholder_pattern = None
        placeholder_config = placeholders or {}
        if placeholder_config:
            escaped = [re.escape(token) for token in placeholder_config.keys()]
            placeholder_pattern = re.compile("|".join(escaped))
        for idx, chunk in enumerate(chunks):
            if idx == 0:
                p = frame.paragraphs[0]
                p.text = ""
            else:
                p = frame.add_paragraph()
            runs = self._build_runs(chunk, default_color, placeholder_config, placeholder_pattern)
            if not runs:
                runs = [(chunk, default_color)]
            for text_part, color in runs:
                if not text_part:
                    continue
                run = p.add_run()
                run.text = text_part
                run.font.name = self.font_family
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(*color)
            if spacing_before_pt is not None:
                p.space_before = Pt(spacing_before_pt)
            else:
                p.space_before = Pt(0)
            if spacing_after_pt is not None:
                p.space_after = Pt(spacing_after_pt)
            else:
                p.space_after = Pt(4)

    def _append_notes(self, textbox, notes: Optional[List[Any]]):
        if not notes:
            return
        frame = textbox.text_frame
        for raw in notes:
            if isinstance(raw, dict):
                text = raw.get("text", "")
                color = tuple(raw.get("color", (198, 38, 38)))
                font_size = raw.get("font_size", 10)
                space_before = raw.get("space_before", 6)
                space_after = raw.get("space_after", 0)
                bold = raw.get("bold", False)
            else:
                text = str(raw)
                color = (198, 38, 38)
                font_size = 10
                space_before = 6
                space_after = 0
                bold = False
            if not text:
                continue
            if len(color) != 3:
                color = (198, 38, 38)
            p = frame.add_paragraph()
            p.text = text
            p.font.name = self.font_family
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*color)
            p.font.bold = bold
            p.space_before = Pt(space_before)
            p.space_after = Pt(space_after)

    def _split_into_paragraphs(self, text: str, count: int) -> List[str]:
        words = text.split()
        if count <= 1 or len(words) <= count:
            return [text]
        chunk_size = len(words) // count
        paragraphs = []
        for i in range(count - 1):
            paragraphs.append(" ".join(words[i * chunk_size:(i + 1) * chunk_size]))
        paragraphs.append(" ".join(words[(count - 1) * chunk_size:]))
        return paragraphs

    def _render_component(self, slide, media_cfg):
        file_path = media_cfg.get("file")
        if not file_path or not Path(file_path).exists():
            print(f"  [WARN] Component file missing: {file_path}")
            return
        class_name = media_cfg.get("class")
        method_name = media_cfg.get("method", "add_to_slide")

        # 如果是流程圖類 component，且 DISABLE_FLOWS=1，就直接跳過
        if self.disable_flows and ("flow" in str(class_name).lower()):
            print(f"  [INFO] Flow component skipped due to DISABLE_FLOWS=1: {class_name}")
            return

        spec = importlib.util.spec_from_file_location("ppt_component", file_path)
        module = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(module)  # type: ignore
        component_class = getattr(module, class_name)
        instance = component_class(self.prs, **media_cfg.get("init_kwargs", {}))
        method = getattr(instance, method_name)
        method(slide, **media_cfg.get("method_kwargs", {}))

    def _draw_overlay(self, slide, overlay_cfg: Dict[str, float]):
        rect = slide.shapes.add_shape(
            1,
            Inches(cm(overlay_cfg.get("left_cm", 0))),
            Inches(cm(overlay_cfg.get("top_cm", 0))),
            Inches(cm(overlay_cfg.get("width_cm", 0))),
            Inches(cm(overlay_cfg.get("height_cm", 0))),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
        rect.fill.transparency = 0.35
        rect.line.color.rgb = RGBColor(200, 200, 200)
        rect.line.width = Pt(0.5)

    def _cleanup_slide(self, slide):
        removable = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                placeholder_format = getattr(shape, "placeholder_format", None)
                if placeholder_format and getattr(placeholder_format, "type", None) == 1:
                    continue
            removable.append(shape)
        for shape in removable:
            slide.shapes._spTree.remove(shape._element)

    def _add_watermark(self, slide, text: Optional[str]):
        if not text:
            return
        width = self.prs.slide_width
        height = self.prs.slide_height
        box_width = Inches(cm(6.0))
        box_height = Inches(cm(1.2))
        left = width - box_width - Inches(cm(1.0))
        top = height - box_height - Inches(cm(1.0))
        textbox = slide.shapes.add_textbox(left, top, box_width, box_height)
        frame = textbox.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = self.font_family
        paragraph.font.size = Pt(11)
        paragraph.font.italic = True
        paragraph.font.color.rgb = RGBColor(150, 150, 150)
        paragraph.alignment = PP_ALIGN.RIGHT

    def _apply_alert(self, shape, text: str):
        if not text:
            return
        try:
            line = shape.line
            line.color.rgb = RGBColor(198, 38, 38)
            line.width = Pt(2)
            line.dash_style = MSO_LINE_DASH_STYLE.SYS_DASH
        except AttributeError:
            return

        slide = shape.part.slide
        label_width = shape.width
        label_height = Pt(16)
        left = shape.left
        top = max(shape.top - label_height - Pt(2), Pt(4))
        textbox = slide.shapes.add_textbox(left, top, label_width, label_height)
        textbox.fill.background()
        textbox.line.fill.background()
        frame = textbox.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = self.font_family
        paragraph.font.size = Pt(9)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(198, 38, 38)
        paragraph.alignment = PP_ALIGN.LEFT

    def _build_runs(
        self,
        text: str,
        default_color: Tuple[int, int, int],
        placeholders: Dict[str, Dict[str, Any]],
        pattern: Optional[re.Pattern],
    ) -> List[Tuple[str, Tuple[int, int, int]]]:
        if not pattern:
            return [(text, default_color)]
        runs: List[Tuple[str, Tuple[int, int, int]]] = []
        pos = 0
        for match in pattern.finditer(text):
            start, end = match.span()
            if start > pos:
                runs.append((text[pos:start], default_color))
            token = match.group()
            config = placeholders.get(token, {})
            replacement = config.get("text", token)
            color = tuple(config.get("color", default_color))
            if len(color) != 3:
                color = default_color
            runs.append((replacement, color))
            pos = end
        if pos < len(text):
            runs.append((text[pos:], default_color))
        return runs


def main():
    print("Generating experimental company + governance + social PPT...")
    engine = PPTFullEngine(PPTContentEngine())
    engine.generate()
    print("Done.")


if __name__ == "__main__":
    main()

