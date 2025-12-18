"""
TCFD 氣候風險時間軸表格 - PPTX 版本（中文）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pathlib import Path


def create_climate_risk_timeline_pptx(output_path=None):
    """建立氣候風險時間軸 PPTX"""
    
    prs = Presentation()
    # 16:9 寬螢幕
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    
    # 標題
    slide_w = 13.333
    title_w = 12.0
    left_margin = (slide_w - title_w) / 2
    
    title = slide.shapes.add_textbox(Inches(left_margin), Inches(0.2), Inches(title_w), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = "氣候相關風險時間軸"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(47, 82, 51)  # 深綠色
    p.alignment = PP_ALIGN.CENTER
    
    # 表格資料（中文版）
    risks_data = [
        ('轉型\n風險', '1', '溫室氣體排放成本增加 - 碳定價機制', '2025', '短期'),
        ('', '2-1', '溫室氣體排放成本增加 - 碳邊境調整(CCA)', '2025', '短期'),
        ('', '2-2', '溫室氣體排放成本增加 - 碳邊境調整(CBAM)', '2026', '中期'),
        ('', '3-1', '永續要求與法規增加 - 其他產業', '2030', '中期'),
        ('', '4', '高碳資產投資風險增加', '2030', '中期'),
        ('', '5', '低碳經濟轉型成本', '2028', '中期'),
        ('', '6-1', '原物料成本增加（基礎建設材料）', '2027', '中期'),
        ('', '3-2', '永續要求與法規增加 - 汽車產業', '2035', '長期'),
        ('', '6-2', '原物料成本增加（食品）', '2050', '長期'),
        ('實體\n風險', '7', '極端氣候事件增加 - 颱風', '2050', '長期'),
        ('', '8', '極端氣候事件增加 - 豪雨', '2050', '長期'),
        ('', '9', '極端氣候事件增加 - 乾旱', '2050', '長期'),
        ('', '10', '平均溫度上升', '2050', '長期'),
        ('', '11', '海平面上升', '2050', '長期')
    ]
    
    # 表格尺寸
    table_w = 12.0
    table_left = (slide_w - table_w) / 2
    rows = len(risks_data) + 1  # +1 for header
    cols = 4
    
    table_shape = slide.shapes.add_table(rows, cols, Inches(table_left), Inches(0.9), Inches(table_w), Inches(6.2))
    tbl = table_shape.table
    
    # 欄寬
    tbl.columns[0].width = Inches(1.2)   # 類別
    tbl.columns[1].width = Inches(0.8)   # 編號
    tbl.columns[2].width = Inches(6.5)   # 風險事件
    tbl.columns[3].width = Inches(3.5)   # 時間軸
    
    # 表頭
    def set_cell(cell, text, bg_color, font_color=RGBColor(255,255,255), bold=True, size=12):
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 合併表頭
    tbl.cell(0, 0).merge(tbl.cell(0, 1))
    set_cell(tbl.cell(0, 0), "風險事件", RGBColor(74, 124, 89))
    
    tbl.cell(0, 2).merge(tbl.cell(0, 3))
    set_cell(tbl.cell(0, 2), "時間軸", RGBColor(107, 148, 117))
    
    # 填入資料
    transition_start = None
    physical_start = None
    
    for i, (category, number, risk_event, year, period) in enumerate(risks_data):
        row_idx = i + 1
        
        # 類別欄
        if category:
            if '轉型' in category:
                transition_start = row_idx
                bg = RGBColor(139, 157, 131)  # 淺綠色
            else:
                physical_start = row_idx
                bg = RGBColor(155, 168, 158)  # 淺灰綠
            
            cell = tbl.cell(row_idx, 0)
            cell.text = category
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 編號欄
        cell = tbl.cell(row_idx, 1)
        cell.text = number
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 風險事件欄
        cell = tbl.cell(row_idx, 2)
        cell.text = risk_event
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 時間軸欄
        cell = tbl.cell(row_idx, 3)
        cell.text = f"{year}    {period}"
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 時期背景色
        if '短期' in period:
            bg = RGBColor(232, 245, 232)  # 很淺綠
        elif '中期' in period:
            bg = RGBColor(212, 230, 212)  # 淺綠
        else:  # 長期
            bg = RGBColor(192, 214, 192)  # 中淺綠
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    
    # 合併類別欄
    # 轉型風險: row 1-9
    tbl.cell(1, 0).merge(tbl.cell(9, 0))
    # 實體風險: row 10-14
    tbl.cell(10, 0).merge(tbl.cell(14, 0))
    
    # 儲存
    if output_path is None:
        output_path = Path(__file__).parent.parent / "output" / "TCFD_氣候風險時間軸.pptx"
    
    prs.save(str(output_path))
    print(f"✓ 氣候風險時間軸表格已生成：{output_path}")
    return output_path


# 給 environment_pptx.py 呼叫的函數（右半邊版本，A4 橫向適配）
def create_tcfd_main_slide_right(slide):
    """在投影片右半邊建立氣候風險時間軸表格（A4 橫向，縮減 2cm）"""
    
    # A4 橫向寬度
    slide_w = 11.69
    
    risks_data = [
        ('轉型\n風險', '1', '溫室氣體排放成本增加 - 碳定價機制', '2025', '短期'),
        ('', '2-1', '溫室氣體排放成本增加 - 碳邊境調整(CCA)', '2025', '短期'),
        ('', '2-2', '溫室氣體排放成本增加 - 碳邊境調整(CBAM)', '2026', '中期'),
        ('', '3-1', '永續要求與法規增加 - 其他產業', '2030', '中期'),
        ('', '4', '高碳資產投資風險增加', '2030', '中期'),
        ('', '5', '低碳經濟轉型成本', '2028', '中期'),
        ('', '6-1', '原物料成本增加（基礎建設材料）', '2027', '中期'),
        ('', '3-2', '永續要求與法規增加 - 汽車產業', '2035', '長期'),
        ('', '6-2', '原物料成本增加（食品）', '2050', '長期'),
        ('實體\n風險', '7', '極端氣候事件增加 - 颱風', '2050', '長期'),
        ('', '8', '極端氣候事件增加 - 豪雨', '2050', '長期'),
        ('', '9', '極端氣候事件增加 - 乾旱', '2050', '長期'),
        ('', '10', '平均溫度上升', '2050', '長期'),
        ('', '11', '海平面上升', '2050', '長期')
    ]
    
    # A4 橫向適配：縮減約 2cm，表格寬度 5.5"
    table_w = 5.5
    # 靠右對齊：slide_w - table_w - 右邊界(0.4)
    table_left = slide_w - table_w - 0.4
    
    # 表格標題（在表格上方）
    title_box = slide.shapes.add_textbox(Inches(table_left), Inches(0.85), Inches(table_w), Inches(0.3))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "氣候相關風險時間軸"
    title_para.font.size = Pt(12)
    title_para.font.bold = True
    title_para.font.name = 'Microsoft JhengHei'
    title_para.font.color.rgb = RGBColor(47, 82, 51)  # 深綠色
    title_para.alignment = PP_ALIGN.CENTER
    
    rows = len(risks_data) + 1
    cols = 4
    
    # 表格位置：標題下方
    table_shape = slide.shapes.add_table(rows, cols, Inches(table_left), Inches(1.2), Inches(table_w), Inches(5.8))
    tbl = table_shape.table
    
    # 欄寬（等比縮小）
    tbl.columns[0].width = Inches(0.55)  # 類別
    tbl.columns[1].width = Inches(0.4)   # 編號
    tbl.columns[2].width = Inches(3.0)   # 風險事件
    tbl.columns[3].width = Inches(1.55)  # 時間軸
    
    # 表頭
    tbl.cell(0, 0).merge(tbl.cell(0, 1))
    _set_header_cell(tbl.cell(0, 0), "風險事件", RGBColor(74, 124, 89), size=10)
    
    tbl.cell(0, 2).merge(tbl.cell(0, 3))
    _set_header_cell(tbl.cell(0, 2), "時間軸", RGBColor(107, 148, 117), size=10)
    
    # 填入資料（斑馬條紋：奇數行淡灰，偶數行白色）
    LIGHT_GRAY = RGBColor(245, 245, 245)  # 淡灰色
    WHITE = RGBColor(255, 255, 255)       # 白色
    
    for i, (category, number, risk_event, year, period) in enumerate(risks_data):
        row_idx = i + 1
        
        # 斑馬條紋：奇數行淡灰，偶數行白色
        row_bg = LIGHT_GRAY if row_idx % 2 == 1 else WHITE
        
        # 左側類別欄（轉型風險/實體風險）
        if category:
            # 統一使用深綠色背景，白色字體
            category_bg = RGBColor(74, 124, 89)  # 深綠色
            _set_category_cell(tbl.cell(row_idx, 0), category, category_bg, size=9)
        else:
            # 非類別欄也設定背景色（斑馬條紋）
            cell = tbl.cell(row_idx, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
        
        # 編號欄（斑馬條紋）
        _set_number_cell_with_bg(tbl.cell(row_idx, 1), number, row_bg, size=8)
        
        # 風險事件欄（斑馬條紋）
        _set_event_cell_with_bg(tbl.cell(row_idx, 2), risk_event, row_bg, size=8)
        
        # 時間軸欄（斑馬條紋，但保留時期顏色）
        _set_timeline_cell_with_bg(tbl.cell(row_idx, 3), year, period, row_bg, size=8)
    
    # 合併類別欄
    tbl.cell(1, 0).merge(tbl.cell(9, 0))
    tbl.cell(10, 0).merge(tbl.cell(14, 0))
    
    print("  ✓ 氣候風險時間軸表格已插入（右半邊）")


# 全頁版本（保留）
def create_tcfd_main_slide(slide):
    """在現有投影片上建立氣候風險時間軸表格（全頁版）"""
    
    slide_w = 13.333
    
    risks_data = [
        ('轉型\n風險', '1', '溫室氣體排放成本增加 - 碳定價機制', '2025', '短期'),
        ('', '2-1', '溫室氣體排放成本增加 - 碳邊境調整(CCA)', '2025', '短期'),
        ('', '2-2', '溫室氣體排放成本增加 - 碳邊境調整(CBAM)', '2026', '中期'),
        ('', '3-1', '永續要求與法規增加 - 其他產業', '2030', '中期'),
        ('', '4', '高碳資產投資風險增加', '2030', '中期'),
        ('', '5', '低碳經濟轉型成本', '2028', '中期'),
        ('', '6-1', '原物料成本增加（基礎建設材料）', '2027', '中期'),
        ('', '3-2', '永續要求與法規增加 - 汽車產業', '2035', '長期'),
        ('', '6-2', '原物料成本增加（食品）', '2050', '長期'),
        ('實體\n風險', '7', '極端氣候事件增加 - 颱風', '2050', '長期'),
        ('', '8', '極端氣候事件增加 - 豪雨', '2050', '長期'),
        ('', '9', '極端氣候事件增加 - 乾旱', '2050', '長期'),
        ('', '10', '平均溫度上升', '2050', '長期'),
        ('', '11', '海平面上升', '2050', '長期')
    ]
    
    table_w = 12.0
    table_left = (slide_w - table_w) / 2
    rows = len(risks_data) + 1
    cols = 4
    
    table_shape = slide.shapes.add_table(rows, cols, Inches(table_left), Inches(1.2), Inches(table_w), Inches(5.8))
    tbl = table_shape.table
    
    # 欄寬
    tbl.columns[0].width = Inches(1.2)
    tbl.columns[1].width = Inches(0.8)
    tbl.columns[2].width = Inches(6.5)
    tbl.columns[3].width = Inches(3.5)
    
    # 表頭
    tbl.cell(0, 0).merge(tbl.cell(0, 1))
    _set_header_cell(tbl.cell(0, 0), "風險事件", RGBColor(74, 124, 89))
    
    tbl.cell(0, 2).merge(tbl.cell(0, 3))
    _set_header_cell(tbl.cell(0, 2), "時間軸", RGBColor(107, 148, 117))
    
    # 填入資料
    for i, (category, number, risk_event, year, period) in enumerate(risks_data):
        row_idx = i + 1
        
        if category:
            bg = RGBColor(139, 157, 131) if '轉型' in category else RGBColor(155, 168, 158)
            _set_category_cell(tbl.cell(row_idx, 0), category, bg)
        
        _set_number_cell(tbl.cell(row_idx, 1), number)
        _set_event_cell(tbl.cell(row_idx, 2), risk_event)
        _set_timeline_cell(tbl.cell(row_idx, 3), year, period)
    
    # 合併類別欄
    tbl.cell(1, 0).merge(tbl.cell(9, 0))
    tbl.cell(10, 0).merge(tbl.cell(14, 0))
    
    print("  ✓ 氣候風險時間軸表格已插入")


def _set_header_cell(cell, text, bg_color, size=12):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_category_cell(cell, text, bg_color, size=9):
    """類別欄（轉型風險/實體風險）- 統一格式"""
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # 白色字體
    p.font.name = 'Microsoft JhengHei'  # 統一字體
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_number_cell(cell, text, size=10):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_event_cell(cell, text, size=10):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.alignment = PP_ALIGN.LEFT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_timeline_cell(cell, year, period, size=10):
    cell.text = f"{year}    {period}"
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if '短期' in period:
        bg = RGBColor(232, 245, 232)
    elif '中期' in period:
        bg = RGBColor(212, 230, 212)
    else:
        bg = RGBColor(192, 214, 192)
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg


# 帶背景色的版本（用於斑馬條紋）
def _set_number_cell_with_bg(cell, text, bg_color, size=8):
    """編號欄（帶背景色）"""
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_event_cell_with_bg(cell, text, bg_color, size=8):
    """風險事件欄（帶背景色）"""
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.alignment = PP_ALIGN.LEFT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_timeline_cell_with_bg(cell, year, period, row_bg, size=8):
    """時間軸欄（帶背景色，但保留時期顏色）"""
    cell.text = f"{year}    {period}"
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 時期顏色（覆蓋斑馬條紋）
    if '短期' in period:
        period_bg = RGBColor(232, 245, 232)  # 很淺綠
    elif '中期' in period:
        period_bg = RGBColor(212, 230, 212)   # 淺綠
    else:  # 長期
        period_bg = RGBColor(192, 214, 192)  # 中淺綠
    
    cell.fill.solid()
    cell.fill.fore_color.rgb = period_bg


# ============ 測試 ============
if __name__ == "__main__":
    create_climate_risk_timeline_pptx()

