"""
Emission 引擎 - PPTX 版本
輸出：
1. 排放表格 PPTX
2. Pie Chart 圖片
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # 使用非互動式後端

# 輸出目錄
OUTPUT_DIR = Path(__file__).parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# ============ 預設排放數據 ============
DEFAULT_EMISSION_DATA = {
    "data_year": "2024",
    "unit": "tCO₂e",
    "scope1": {
        "gasoline": 1.16,
        "refrigerant": 4.18,
        "subtotal": 5.34
    },
    "scope2": {
        "electricity": 148.11,
        "subtotal": 148.11
    },
    "scope3": {
        "goods_services": 0.00,
        "transportation": 0.00,
        "subtotal": 0.00
    },
    "total": 153.45
}

# 動態數據（可由外部設定）
EMISSION_DATA = DEFAULT_EMISSION_DATA.copy()

def set_emission_data(data):
    """設定動態排放數據（從 Emission Engine 傳入）"""
    global EMISSION_DATA
    if data:
        # 取得細項數據（優先使用 gasoline/refrigerant，否則用 scope1 均分）
        s1_total = data.get("scope1", 0)
        s1_gasoline = data.get("gasoline", s1_total * 0.5)
        s1_refrigerant = data.get("refrigerant", s1_total * 0.5)
        
        s2_total = data.get("scope2", 0)
        s2_electricity = data.get("electricity", s2_total)
        
        total = data.get("total", s1_total + s2_total)
        
        EMISSION_DATA = {
            "data_year": data.get("data_year", "2024"),
            "unit": "tCO₂e",
            "scope1": {
                "gasoline": s1_gasoline,
                "refrigerant": s1_refrigerant,
                "subtotal": s1_total
            },
            "scope2": {
                "electricity": s2_electricity,
                "subtotal": s2_total
            },
            "scope3": {
                "goods_services": 0.00,
                "transportation": 0.00,
                "subtotal": data.get("scope3", 0)
            },
            "total": total
        }
        print(f"  ✓ 已載入動態排放數據：")
        print(f"    範疇一: {s1_total:.2f} t (汽油 {s1_gasoline:.2f}, 冷媒 {s1_refrigerant:.2f})")
        print(f"    範疇二: {s2_total:.2f} t")
        print(f"    總排放: {total:.2f} tCO₂e")


def set_cell_fill(cell, hex_color):
    """設定儲存格背景顏色"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = OxmlElement('a:solidFill')
    srgbClr = OxmlElement('a:srgbClr')
    srgbClr.set('val', hex_color)
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def create_emission_table_pptx(output_path=None):
    """建立排放表格 PPTX（A4 橫向）"""
    prs = Presentation()
    # A4 橫向: 297mm x 210mm = 11.69" x 8.27"
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "溫室氣體排放盤查表 (範疇 1-3)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = 'Microsoft JhengHei'
    p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), Inches(12), Inches(0.4))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"資料年度: {EMISSION_DATA['data_year']} | 單位: {EMISSION_DATA['unit']}"
    p2.font.size = Pt(14)
    p2.font.name = 'Microsoft JhengHei'
    p2.alignment = PP_ALIGN.CENTER
    
    # 建立表格
    rows = 10
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.4), Inches(10), Inches(5)).table
    
    # 設定欄寬
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.5)
    
    # 表頭
    headers = ["範疇", "排放來源", "排放量\n(tCO₂e)", "備註"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        set_cell_fill(cell, "2F5233")  # 深綠色
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.name = 'Microsoft JhengHei'
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 資料列（使用動態數據）
    s1_gas = EMISSION_DATA["scope1"]["gasoline"]
    s1_ref = EMISSION_DATA["scope1"]["refrigerant"]
    s1_sub = EMISSION_DATA["scope1"]["subtotal"]
    s2_elec = EMISSION_DATA["scope2"]["electricity"]
    s2_sub = EMISSION_DATA["scope2"]["subtotal"]
    s3_sub = EMISSION_DATA["scope3"]["subtotal"]
    total = EMISSION_DATA["total"]
    
    # 計算占比
    s2_percent = (s2_sub / total * 100) if total > 0 else 0
    
    data_rows = [
        ["範疇一", "汽油（公務車）", f"{s1_gas:.2f}", "估算"],
        ["", "冷媒（R-410A）", f"{s1_ref:.2f}", "維護估算"],
        ["小計", "", f"{s1_sub:.2f}", ""],
        ["範疇二", "外購電力", f"{s2_elec:.2f}", f"佔總排放{s2_percent:.1f}%"],
        ["小計", "", f"{s2_sub:.2f}", ""],
        ["範疇三", "採購商品與服務", "0.00", "暫未納入"],
        ["", "運輸配送", "0.00", "暫未納入"],
        ["小計", "", f"{s3_sub:.2f}", ""],
        ["合計", "", f"{total:.2f}", "100%"],
    ]
    
    for row_idx, row_data in enumerate(data_rows, start=1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.name = 'Microsoft JhengHei'
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 特殊列樣式
            if row_data[0] in ["小計", "合計"]:
                p.font.bold = True
                if row_data[0] == "合計":
                    set_cell_fill(cell, "E8F5E9")  # 淺綠色
    
    # 儲存
    if output_path is None:
        output_path = OUTPUT_DIR / "emission_table.pptx"
    
    prs.save(str(output_path))
    print(f"✓ 排放表格已儲存: {output_path}")
    
    return output_path


def create_emission_pie_chart(output_path=None):
    """建立排放圓餅圖"""
    # 設定中文字體
    plt.rcParams['font.sans-serif'] = ['Microsoft JhengHei', 'SimHei', 'Arial']
    plt.rcParams['axes.unicode_minus'] = False
    
    scope1 = EMISSION_DATA["scope1"]["subtotal"]
    scope2 = EMISSION_DATA["scope2"]["subtotal"]
    scope3 = EMISSION_DATA["scope3"]["subtotal"]
    
    labels = ["範疇一\n(直接排放)", "範疇二\n(外購電力)", "範疇三\n(其他間接)"]
    sizes = [scope1, scope2, scope3]
    colors = ["#6BA292", "#007A3D", "#C1C1C1"]
    explode = (0, 0.05, 0)  # 突顯範疇二
    
    # 避免 0 值造成問題
    if sum(sizes) == 0:
        sizes = [1, 1, 1]
    
    fig, ax = plt.subplots(figsize=(8, 6), dpi=150)
    
    wedges, texts, autotexts = ax.pie(
        sizes,
        labels=labels,
        autopct=lambda p: f"{p:.1f}%\n({p/100*sum([scope1, scope2, scope3]):.2f} t)" if p > 0 else "",
        startangle=90,
        colors=colors,
        explode=explode,
        textprops={"fontsize": 11, "fontweight": "bold"}
    )
    
    ax.set_title("溫室氣體排放佔比 (tCO₂e)", fontsize=16, fontweight='bold', pad=20)
    ax.axis('equal')
    
    # 圖例
    ax.legend(wedges, [f"範疇一: {scope1:.2f} t", f"範疇二: {scope2:.2f} t", f"範疇三: {scope3:.2f} t"],
              title="排放量", loc="center left", bbox_to_anchor=(1, 0, 0.5, 1))
    
    plt.tight_layout()
    
    if output_path is None:
        output_path = OUTPUT_DIR / "emission_pie_chart.png"
    
    plt.savefig(str(output_path), bbox_inches="tight", dpi=300, facecolor='white')
    plt.close()
    
    print(f"✓ 圓餅圖已儲存: {output_path}")
    
    return output_path


def create_emission_table_on_slide_right(slide):
    """在投影片右半邊建立排放表格（縮減50%寬度，靠右對齊）"""
    # A4 橫向寬度
    slide_w = 11.69
    
    # 原始表格寬度 10"，縮減50% = 5"
    table_w = 5.0
    # 靠右對齊：slide_w - table_w - 右邊界(0.4)
    table_left = slide_w - table_w - 0.4
    
    # 表格位置：top 調整為 1.2" 以配合左側文字
    rows = 10
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(table_left), Inches(1.2), Inches(table_w), Inches(5.8))
    tbl = table_shape.table
    
    # 欄寬（等比縮小50%）
    tbl.columns[0].width = Inches(1.0)   # 範疇
    tbl.columns[1].width = Inches(1.75) # 排放來源
    tbl.columns[2].width = Inches(1.0)   # 排放量
    tbl.columns[3].width = Inches(1.25) # 備註
    
    # 表頭
    headers = ["範疇", "排放來源", "排放量\n(tCO₂e)", "備註"]
    for i, header in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = header
        set_cell_fill(cell, "2F5233")  # 深綠色
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(9)  # 縮小字體
        p.font.name = 'Microsoft JhengHei'
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 資料列（使用動態數據）
    s1_gas = EMISSION_DATA["scope1"]["gasoline"]
    s1_ref = EMISSION_DATA["scope1"]["refrigerant"]
    s1_sub = EMISSION_DATA["scope1"]["subtotal"]
    s2_elec = EMISSION_DATA["scope2"]["electricity"]
    s2_sub = EMISSION_DATA["scope2"]["subtotal"]
    s3_sub = EMISSION_DATA["scope3"]["subtotal"]
    total = EMISSION_DATA["total"]
    
    # 計算占比
    s2_percent = (s2_sub / total * 100) if total > 0 else 0
    
    data_rows = [
        ["範疇一", "汽油（公務車）", f"{s1_gas:.2f}", "估算"],
        ["", "冷媒（R-410A）", f"{s1_ref:.2f}", "維護估算"],
        ["小計", "", f"{s1_sub:.2f}", ""],
        ["範疇二", "外購電力", f"{s2_elec:.2f}", f"佔總排放{s2_percent:.1f}%"],
        ["小計", "", f"{s2_sub:.2f}", ""],
        ["範疇三", "採購商品與服務", "0.00", "暫未納入"],
        ["", "運輸配送", "0.00", "暫未納入"],
        ["小計", "", f"{s3_sub:.2f}", ""],
        ["合計", "", f"{total:.2f}", "100%"],
    ]
    
    for row_idx, row_data in enumerate(data_rows, start=1):
        for col_idx, value in enumerate(row_data):
            cell = tbl.cell(row_idx, col_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)  # 縮小字體
            p.font.name = 'Microsoft JhengHei'
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 特殊列樣式
            if row_data[0] in ["小計", "合計"]:
                p.font.bold = True
                if row_data[0] == "合計":
                    set_cell_fill(cell, "E8F5E9")  # 淺綠色
    
    print("  ✓ 排放表格已插入（右半邊，縮減50%）")


def generate_all():
    """生成所有排放相關檔案"""
    print("\n[生成排放引擎輸出]")
    table_path = create_emission_table_pptx()
    chart_path = create_emission_pie_chart()
    
    return {
        "table_pptx": table_path,
        "pie_chart": chart_path
    }


if __name__ == "__main__":
    results = generate_all()
    print(f"\n完成！")
    print(f"  表格: {results['table_pptx']}")
    print(f"  圖表: {results['pie_chart']}")

