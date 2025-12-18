"""
ESG 報告生成器 - 環境篇 PPTX 引擎
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os
import glob
import copy
from lxml import etree

from config import ENVIRONMENT_CONFIG, ENVIRONMENT_IMAGE_MAPPING, TCFD_TABLES, ASSETS_PATH
from content_engine import ContentEngine

# 加入 assets 路徑
import sys
sys.path.insert(0, str(Path(__file__).parent / "assets"))
from TCFD_main_pptx import create_tcfd_main_slide_right
from emission_pptx import create_emission_table_on_slide_right
sys.path.append(ASSETS_PATH)

# ============ SASB 產業映射 ============
SASB_MAP = {
    # 采矿冶金
    "钢铁": ("EM-IS", "钢铁生产"), "水泥": ("EM-CM", "建材"), "煤矿": ("EM-CO", "煤炭"),
    "采矿": ("EM-MM", "金属矿业"), "石油": ("EM-OG", "油气开采"),
    
    # 能源公用
    "电力": ("IF-EU", "电力公用事业"), "天然气": ("IF-GU", "燃气"), "自来水": ("IF-WU", "水务"),
    
    # 运输
    "航空": ("TR-AF", "航空"), "物流": ("TR-AL", "空运物流"), "铁路": ("TR-RA", "铁路"),
    "货运": ("TR-RO", "公路运输"), "海运": ("TR-MT", "海运"), "租车": ("TR-CR", "租车"),
    
    # 房地产建筑
    "房地产": ("IF-RE", "不动产"), "建商": ("IF-HB", "建筑商"), "工程": ("IF-EN", "工程建设"),
    
    # 食品饮料
    "农业": ("FB-AB", "农产品"), "畜牧": ("FB-MP", "畜牧乳品"), "食品": ("FB-PF", "食品加工"),
    "餐饮": ("FB-RN", "餐饮"), "超市": ("FB-FR", "食品零售"), "饮料": ("FB-NB", "饮料"),
    
    # 医疗
    "制药": ("HC-BI", "生技制药"), "医材": ("HC-MS", "医疗器材"), "药局": ("HC-DY", "药品零售"),
    "医院": ("HC-DV", "医疗服务"),
    
    # 金融
    "银行": ("FN-CB", "商业银行"), "证券": ("FN-IB", "投资银行"), "保险": ("FN-IN", "保险"),
    "基金": ("FN-AC", "资产管理"), "消金": ("FN-CF", "消费金融"),
    
    # 科技
    "半导体": ("TC-SC", "半导体"), "电子": ("TC-HW", "硬件"), "软件": ("TC-SI", "软件IT"),
    "电信": ("TC-TL", "电信"), "电商": ("TC-ES", "电子商务"), "互联网": ("TC-IM", "网络媒体"),
    
    # 消费品
    "服装": ("CG-AA", "服饰"), "家电": ("CG-AM", "家电"), "家具": ("CG-BF", "家具"),
    "日用品": ("CG-HP", "家用品"), "零售": ("CG-MR", "零售"), "汽车": ("CG-AC", "汽车"),
    
    # 其他
    "化工": ("RT-CH", "化工"), "包装": ("RT-CP", "包装"), "机械": ("RT-IG", "工业机械"),
    "造纸": ("RR-PP", "纸浆造纸"), "太阳能": ("RR-ST", "太阳能"),
}

def get_sasb(industry_input):
    """根據產業名稱取得 SASB 代碼和名稱"""
    for keyword, (code, name) in SASB_MAP.items():
        if keyword in industry_input:
            return code, name
    return "通用", "一般产业"

# ============ A4 橫向版面常數 ============
# A4 橫向: 297mm x 210mm = 11.69" x 8.27"
SLIDE_WIDTH = Inches(11.69)
SLIDE_HEIGHT = Inches(8.27)

# 版面配置常數
MARGIN = Inches(0.4)                    # 頁面邊距
TITLE_TOP = Inches(0.3)                 # 標題頂部位置
TITLE_HEIGHT = Inches(0.7)              # 標題高度
CONTENT_TOP = Inches(1.1)               # 內容區頂部
CONTENT_HEIGHT = Inches(6.5)            # 內容區高度
CONTENT_WIDTH = Inches(5.2)             # 單側內容寬度（左或右）
LEFT_CONTENT_LEFT = MARGIN              # 左側內容起點
RIGHT_CONTENT_LEFT = Inches(6.1)        # 右側內容起點
GAP = Inches(0.3)                       # 左右間距

# TCFD Generator 的輸出路徑
TCFD_OUTPUT_PATH = r"C:\Users\User\Desktop\TCFD generator\output"

# Emission 輸出路徑
EMISSION_OUTPUT_PATH = Path(__file__).parent / "output"


class EnvironmentPPTXEngine:
    """環境篇 PPTX 報告生成引擎"""

    def __init__(self, template_path=None, test_mode=False, emission_data=None, industry="企業", tcfd_output_folder=None, emission_output_folder=None, company_profile=None, api_key=None):
        """
        初始化引擎
        template_path: 模板檔案路徑（可選）
        test_mode: 測試模式（跳過 LLM API）
        emission_data: 碳排放數據 dict
        industry: 產業名稱
        tcfd_output_folder: TCFD 輸出資料夾路徑（從 Step 1 傳入）
        emission_output_folder: Emission 輸出資料夾路徑（從 Step 2 傳入）
        company_profile: 公司規模資訊 dict（從 Step 2 傳入）
        api_key: Claude API Key
        """
        self.emission_data = emission_data or {}
        self.industry = industry
        self.tcfd_output_folder = tcfd_output_folder  # Step 1 的 TCFD 輸出路徑
        self.emission_output_folder = emission_output_folder  # Step 2 的 Emission 輸出路徑
        self.company_profile = company_profile or {}
        self.api_key = api_key
        if template_path and os.path.exists(template_path):
            # 方案 C：備份模板到 assets 資料夾（使用絕對路徑）
            import shutil
            env_report_dir = Path(__file__).parent
            template_backup_path = env_report_dir / "assets" / "template_backup.pptx"
            template_backup_path.parent.mkdir(exist_ok=True)  # 確保資料夾存在
            if not template_backup_path.exists():
                shutil.copy(template_path, template_backup_path)
                print(f"✓ 模板已備份至: {template_backup_path}")
            
            self.prs = Presentation(template_path)
            print(f"✓ 載入模板：{template_path}")
            
            # 刪除模板自帶的空白頁面
            template_slide_count = len(self.prs.slides)
            if template_slide_count > 0:
                print(f"  ℹ 刪除模板預設頁面 ({template_slide_count} 頁)")
                # 從後往前刪除，避免索引問題
                for i in range(template_slide_count - 1, -1, -1):
                    rId = self.prs.slides._sldIdLst[i].rId
                    self.prs.part.drop_rel(rId)
                    del self.prs.slides._sldIdLst[i]
        else:
            self.prs = Presentation()
        
        # 統一設定為 A4 橫向
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        print(f"✓ 投影片尺寸：A4 橫向 ({SLIDE_WIDTH.inches}\" x {SLIDE_HEIGHT.inches}\")")
        
        self.test_mode = test_mode
        
        # 調試：檢查 API Key
        if self.api_key:
            print(f"  ✓ EnvironmentPPTXEngine 收到 API Key: {self.api_key[:10]}...")
        else:
            print(f"  ⚠ EnvironmentPPTXEngine 未收到 API Key")
        
        self.content_engine = ContentEngine(
            test_mode=test_mode, 
            company_profile=self.company_profile,
            api_key=self.api_key
        )
        self.config = ENVIRONMENT_CONFIG
        
        # 版面設定
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        
        # 樣式設定
        self.title_font_size = Pt(28)
        self.body_font_size = Pt(14)
        self.font_name = 'Microsoft JhengHei'
        self.primary_color = RGBColor(26, 58, 46)  # 深綠色
        self.secondary_color = RGBColor(74, 124, 89)  # 淺綠色

    def _get_blank_layout(self):
        """取得空白版面配置"""
        # 優先使用 index 6（標準空白版面）
        if len(self.prs.slide_layouts) > 6:
            return self.prs.slide_layouts[6]
        
        # 嘗試找名稱含 blank 的版面
        for layout in self.prs.slide_layouts:
            if 'blank' in layout.name.lower() or layout.name == '空白':
                return layout
        
        # 最後使用第一個版面
        return self.prs.slide_layouts[0]

    def _add_slide(self):
        """新增空白投影片（移除預設佔位符，並添加浮水印）"""
        layout = self._get_blank_layout()
        slide = self.prs.slides.add_slide(layout)
        
        # 移除版面自帶的佔位符形狀
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        
        # 添加浮水印（右下角）
        self._add_watermark(slide)
        
        return slide
    
    def _add_watermark(self, slide):
        """在投影片右下角添加浮水印"""
        try:
            # 計算右下角位置
            watermark_width = Inches(2.5)
            watermark_height = Inches(0.3)
            watermark_left = SLIDE_WIDTH - watermark_width - MARGIN
            watermark_top = SLIDE_HEIGHT - watermark_height - Inches(0.2)
            
            # 添加文字框
            watermark_box = slide.shapes.add_textbox(watermark_left, watermark_top, watermark_width, watermark_height)
            watermark_frame = watermark_box.text_frame
            watermark_para = watermark_frame.paragraphs[0]
            watermark_para.text = "Sustainability Report"
            
            # 設定格式：斜體、加粗、11pt、淡灰色
            watermark_para.font.size = Pt(11)
            watermark_para.font.bold = True
            watermark_para.font.italic = True
            watermark_para.font.name = 'Arial'
            watermark_para.font.color.rgb = RGBColor(200, 200, 200)  # 淡灰色
            watermark_para.alignment = PP_ALIGN.RIGHT
            
        except Exception as e:
            print(f"  ⚠ 浮水印添加警告：{e}")

    def _add_title(self, slide, title_text, left=None, top=None, 
                   width=None, height=None):
        """在投影片上新增標題（使用 A4 橫向常數）"""
        left = left or MARGIN
        top = top or TITLE_TOP
        width = width or Inches(10.8)  # A4 橫向寬度扣邊距
        height = height or TITLE_HEIGHT
        """在投影片上新增標題"""
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.text = title_text
        p.font.size = self.title_font_size
        p.font.bold = True
        p.font.name = self.font_name
        p.font.color.rgb = self.primary_color
        p.alignment = PP_ALIGN.LEFT
        
        return title_box

    def _add_text_box(self, slide, text, left, top, width, height, 
                      font_size=None, bold=False, alignment=PP_ALIGN.LEFT):
        """在投影片上新增文字框"""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = font_size or self.body_font_size
        p.font.bold = bold
        p.font.name = self.font_name
        p.alignment = alignment
        
        return text_box

    def _add_image(self, slide, image_name, left, top, width=None, height=None):
        """在投影片上新增圖片（從 assets 資料夾）"""
        image_path = os.path.join(ASSETS_PATH, image_name)
        return self._add_image_full_path(slide, image_path, left, top, width, height)

    def _add_image_full_path(self, slide, image_path, left, top, width=None, height=None):
        """在投影片上新增圖片（完整路徑）"""
        if os.path.exists(image_path):
            try:
                if width and height:
                    pic = slide.shapes.add_picture(image_path, left, top, width, height)
                elif width:
                    pic = slide.shapes.add_picture(image_path, left, top, width=width)
                elif height:
                    pic = slide.shapes.add_picture(image_path, left, top, height=height)
                else:
                    pic = slide.shapes.add_picture(image_path, left, top)
                print(f"  ✓ 插入圖片：{os.path.basename(image_path)}")
                return pic
            except Exception as e:
                print(f"  ✗ 圖片插入失敗：{image_path} - {e}")
        else:
            print(f"  ✗ 圖片不存在：{image_path}")
        return None

    def _create_left_image_right_text_slide(self, title, image_name, text_content):
        """建立左圖右文版面（A4 橫向）"""
        slide = self._add_slide()
        
        # 標題
        self._add_title(slide, title)
        
        # 左側圖片
        self._add_image(slide, image_name, 
                       left=LEFT_CONTENT_LEFT, 
                       top=CONTENT_TOP, 
                       width=CONTENT_WIDTH,
                       height=CONTENT_HEIGHT)
        
        # 右側文字
        self._add_text_box(slide, text_content,
                          left=RIGHT_CONTENT_LEFT,
                          top=CONTENT_TOP,
                          width=CONTENT_WIDTH,
                          height=CONTENT_HEIGHT)
        
        return slide

    def _create_left_image_right_text_slide_full_path(self, title, image_path, text_content):
        """建立左圖右文版面（A4 橫向，圖片使用完整路徑）"""
        slide = self._add_slide()
        
        # 標題
        self._add_title(slide, title)
        
        # 左側圖片
        self._add_image_full_path(slide, image_path, 
                                  left=LEFT_CONTENT_LEFT, 
                                  top=CONTENT_TOP, 
                                  width=CONTENT_WIDTH,
                                  height=CONTENT_HEIGHT)
        
        # 右側文字
        self._add_text_box(slide, text_content,
                          left=RIGHT_CONTENT_LEFT,
                          top=CONTENT_TOP,
                          width=CONTENT_WIDTH,
                          height=CONTENT_HEIGHT)
        
        return slide

    def _create_left_text_right_image_slide(self, title, text_content, image_name):
        """建立左文右圖版面（A4 橫向）"""
        slide = self._add_slide()
        
        # 標題
        self._add_title(slide, title)
        
        # 左側文字
        self._add_text_box(slide, text_content,
                          left=LEFT_CONTENT_LEFT,
                          top=CONTENT_TOP,
                          width=CONTENT_WIDTH,
                          height=CONTENT_HEIGHT)
        
        # 右側圖片
        self._add_image(slide, image_name,
                       left=RIGHT_CONTENT_LEFT,
                       top=CONTENT_TOP,
                       width=CONTENT_WIDTH,
                       height=CONTENT_HEIGHT)
        
        return slide

    def _create_full_text_slide(self, title, text_content):
        """建立純文字版面（A4 橫向）"""
        slide = self._add_slide()
        
        # 標題
        self._add_title(slide, title)
        
        # 全版文字
        self._add_text_box(slide, text_content,
                          left=MARGIN,
                          top=CONTENT_TOP,
                          width=Inches(10.8),  # A4 橫向寬度扣邊距
                          height=CONTENT_HEIGHT)
        
        return slide

    # ==================== 各章節生成方法 ====================

    def generate_cover_page(self):
        """生成封面頁"""
        print("\n[生成封面頁]")
        
        cover_text = self.content_engine.generate_environmental_cover(self.config)
        self._create_left_image_right_text_slide(
            "第四章 環境永續",
            ENVIRONMENT_IMAGE_MAPPING['cover'],
            cover_text
        )
        
        print("✓ 封面頁完成")

    def generate_policy_pages(self):
        """生成 4.1-4.2 環境政策頁面"""
        print("\n[生成環境政策頁面]")
        
        # Page 1: 4.1 環境政策與管理架構
        sustainability_text = self.content_engine.generate_sustainability_committee(self.config)
        self._create_left_image_right_text_slide(
            "4.1 環境政策與管理架構",
            ENVIRONMENT_IMAGE_MAPPING['sustainability_panel'],
            sustainability_text
        )
        
        # Page 2: 4.2 環境政策四大面向
        policy_text = self.content_engine.generate_policy_description(self.config)
        self._create_left_text_right_image_slide(
            "4.2 環境政策四大面向",
            policy_text,
            ENVIRONMENT_IMAGE_MAPPING['policy']
        )
        
        print("✓ 環境政策頁面完成")

    def _find_latest_tcfd_files(self):
        """尋找 5 個 TCFD PPTX 檔案（優先使用傳入的資料夾）"""
        
        tcfd_patterns = [
            ("01", "TCFD_01_轉型風險_*.pptx"),
            ("02", "TCFD_02_市場風險_*.pptx"),
            ("03", "TCFD_03_實體風險_*.pptx"),
            ("04", "TCFD_04_溫升風險_*.pptx"),
            ("05", "TCFD_05_資源效率_*.pptx"),
        ]
        
        tcfd_files = {}
        
        # 優先使用傳入的 TCFD 資料夾（從 Step 1）
        if self.tcfd_output_folder and os.path.exists(self.tcfd_output_folder):
            search_base = self.tcfd_output_folder
            print(f"  ✓ 使用 Step 1 輸出資料夾: {search_base}")
        else:
            # 回退：尋找最新的子資料夾
            tcfd_folders = glob.glob(os.path.join(TCFD_OUTPUT_PATH, "TCFD_*"))
            tcfd_folders = [f for f in tcfd_folders if os.path.isdir(f)]
            
            if tcfd_folders:
                latest_folder = max(tcfd_folders, key=os.path.getmtime)
                print(f"  ✓ 找到最新 TCFD 資料夾: {os.path.basename(latest_folder)}")
                search_base = latest_folder
            else:
                print(f"  ℹ 從 output 資料夾直接搜尋")
                search_base = TCFD_OUTPUT_PATH
        
        for key, pattern in tcfd_patterns:
            search_path = os.path.join(search_base, pattern)
            files = glob.glob(search_path)
            if files:
                # 取最新的檔案
                latest_file = max(files, key=os.path.getmtime)
                tcfd_files[key] = latest_file
                print(f"  ✓ 找到 TCFD {key}: {os.path.basename(latest_file)}")
            else:
                print(f"  ✗ 未找到 TCFD {key} 檔案")
        
        return tcfd_files

    def _copy_slide_from_pptx(self, source_pptx_path):
        """從另一個 PPTX 複製投影片內容"""
        try:
            source_prs = Presentation(source_pptx_path)
            
            for source_slide in source_prs.slides:
                # 新增空白投影片
                new_slide = self._add_slide()
                
                # 複製所有形狀
                for shape in source_slide.shapes:
                    try:
                        # 複製形狀的 XML
                        el = shape.element
                        new_el = copy.deepcopy(el)
                        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                    except Exception as e:
                        print(f"    ⚠ 複製形狀時警告: {e}")
                
            return True
        except Exception as e:
            print(f"  ✗ 複製投影片失敗: {e}")
            return False

    def _insert_tcfd_pptx(self, pptx_path, title):
        """插入 TCFD PPTX 檔案的投影片"""
        try:
            source_prs = Presentation(pptx_path)
            
            for idx, source_slide in enumerate(source_prs.slides):
                # 新增空白投影片（使用模板的白色背景）
                new_slide = self._add_slide()
                
                # 不複製背景，保持新投影片的白色背景
                # 明確設定背景為白色（確保不會是黑色）
                try:
                    new_slide.background.fill.solid()
                    new_slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色
                except:
                    pass
                
                # 複製所有形狀（但不包括背景）
                for shape in source_slide.shapes:
                    try:
                        el = shape.element
                        new_el = copy.deepcopy(el)
                        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                    except Exception as e:
                        # 如果複製失敗，嘗試用其他方式
                        pass
                
            print(f"  ✓ 插入 {title}: {os.path.basename(pptx_path)}")
            return True
        except Exception as e:
            print(f"  ✗ 插入失敗 {title}: {e}")
            return False

    def generate_tcfd_pages(self):
        """生成 TCFD 頁面（從 TCFD Generator 插入）"""
        print("\n[生成 TCFD 頁面]")
        
        # 4.3 TCFD 氣候財務揭露（左文右表）
        section_slide = self._add_slide()
        self._add_title(section_slide, "4.3 TCFD 氣候財務揭露")
        
        # 左邊：文字說明（使用專用 prompt）
        tcfd_intro_text = self.content_engine.generate_tcfd_financial_disclosure(self.config)
        self._add_text_box(section_slide, tcfd_intro_text,
                          left=LEFT_CONTENT_LEFT, 
                          top=CONTENT_TOP, 
                          width=CONTENT_WIDTH, 
                          height=CONTENT_HEIGHT,
                          font_size=Pt(12))
        
        # 右邊：氣候風險時間軸表格（縮減40%，靠右對齊）
        create_tcfd_main_slide_right(section_slide)
        
        # 尋找 TCFD PPTX 檔案
        tcfd_files = self._find_latest_tcfd_files()
        
        tcfd_info = [
            ("01", "轉型風險分析"),
            ("02", "市場風險分析"),
            ("03", "實體風險分析"),
            ("04", "溫升風險分析"),
            ("05", "資源效率機會"),
        ]
        
        for key, title in tcfd_info:
            if key in tcfd_files:
                # 插入 TCFD PPTX
                self._insert_tcfd_pptx(tcfd_files[key], title)
            else:
                # 如果找不到檔案，建立佔位頁面
                slide = self._add_slide()
                self._add_title(slide, title)
                self._add_text_box(slide, 
                                 f"【{title}】\n\n請先執行 TCFD Generator 生成表格",
                                 left=Inches(0.5),
                                 top=Inches(1.5),
                                 width=Inches(12),
                                 height=Inches(5))
        
        # 4.4 TCFD 氣候風險（風險矩陣）
        matrix_text = self.content_engine.generate_tcfd_matrix_analysis(self.config)
        self._create_left_text_right_image_slide(
            "4.4 TCFD 氣候風險",
            matrix_text,
            ENVIRONMENT_IMAGE_MAPPING['tcfd_matrix']
        )
        
        # Page 10: SASB 產業分類（左文右表，插入在 4.4 和 4.5 之間）
        sasb_slide = self._add_slide()
        self._add_title(sasb_slide, "SASB 產業分類")
        
        # 取得 SASB 代碼和名稱
        sasb_code, sasb_name = get_sasb(self.industry)
        
        # 左邊：LLM 生成的分析文字（280字）
        sasb_analysis_text = self.content_engine.generate_sasb_analysis(
            self.config, 
            self.industry, 
            sasb_code, 
            sasb_name
        )
        self._add_text_box(sasb_slide, sasb_analysis_text,
                          left=LEFT_CONTENT_LEFT, 
                          top=CONTENT_TOP, 
                          width=CONTENT_WIDTH, 
                          height=CONTENT_HEIGHT,
                          font_size=Pt(12))
        
        # 右邊：SASB 表格（A4 一半寬度，靠右對齊）
        self._create_sasb_table(sasb_slide, sasb_code, sasb_name)
        
        print("✓ TCFD 頁面完成")

    def _create_sasb_table(self, slide, sasb_code, sasb_name):
        """在投影片右半邊建立 SASB 表格（寬度減少 1.5 公分，列出所有產業分類）"""
        from pptx.oxml.xmlchemy import OxmlElement
        
        # A4 橫向寬度
        slide_w = 11.69
        
        # 表格寬度：A4 一半 - 1.5 公分 = 5.845" - 0.59" = 5.255"
        table_w = (slide_w / 2) - 0.59
        # 靠右對齊：slide_w - table_w - 右邊界(0.4)
        table_left = slide_w - table_w - 0.4
        
        # SASB 產業大類列表
        sasb_categories = [
            ("CG", "Consumer Goods", "消费品"),
            ("EM", "Extractives & Minerals Processing", "採掘與礦物加工"),
            ("FN", "Financials", "金融服務"),
            ("FB", "Food & Beverage", "食品飲料"),
            ("HC", "Health Care", "醫療保健"),
            ("IF", "Infrastructure", "基礎設施"),
            ("RR", "Renewable Resources & Alternative Energy", "可再生資源與替代能源"),
            ("RT", "Resource Transformation", "資源轉換（製造業）"),
            ("SV", "Services", "服務業"),
            ("TC", "Technology & Communications", "科技與通訊"),
            ("TR", "Transportation", "運輸"),
        ]
        
        # 取得當前產業的大類代碼（前綴）
        current_category_prefix = sasb_code.split("-")[0] if "-" in sasb_code else sasb_code
        
        # 表格位置：top 調整為 1.2" 以配合左側文字
        rows = 1 + len(sasb_categories)  # 表頭 + 11 個分類
        cols = 3  # 代碼、英文名稱、中文名稱
        table_height = min(Inches(5.5), Inches(0.4 * rows))  # 動態高度，最多 5.5"
        table_shape = slide.shapes.add_table(rows, cols, Inches(table_left), Inches(1.2), Inches(table_w), table_height)
        tbl = table_shape.table
        
        # 欄寬
        tbl.columns[0].width = Inches(0.8)   # 代碼
        tbl.columns[1].width = Inches(2.2)   # 英文名稱
        tbl.columns[2].width = Inches(2.255) # 中文名稱
        
        # 設定儲存格背景顏色
        def set_cell_fill(cell, hex_color):
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = OxmlElement('a:solidFill')
            srgbClr = OxmlElement('a:srgbClr')
            srgbClr.set('val', hex_color)
            solidFill.append(srgbClr)
            tcPr.append(solidFill)
        
        # 表頭
        header_cell = tbl.cell(0, 0)
        header_cell.merge(tbl.cell(0, 2))
        header_cell.text = "SASB 產業分類"
        set_cell_fill(header_cell, "2F5233")  # 深綠色
        p = header_cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Microsoft JhengHei'
        p.alignment = PP_ALIGN.CENTER
        header_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 填入所有產業分類
        for row_idx, (code, en_name, zh_name) in enumerate(sasb_categories, start=1):
            # 判斷是否為當前產業的分類
            is_current = (code == current_category_prefix)
            
            # 代碼欄
            code_cell = tbl.cell(row_idx, 0)
            code_cell.text = code
            if is_current:
                set_cell_fill(code_cell, "4A7C59")  # 深綠色（當前分類）
            else:
                set_cell_fill(code_cell, "F1F8F4")  # 極淺綠色
            p = code_cell.text_frame.paragraphs[0]
            if is_current:
                p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(9)
            p.font.bold = is_current
            p.font.name = 'Microsoft JhengHei'
            p.alignment = PP_ALIGN.CENTER
            code_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 英文名稱欄
            en_cell = tbl.cell(row_idx, 1)
            en_cell.text = en_name
            if is_current:
                set_cell_fill(en_cell, "4A7C59")  # 深綠色（當前分類）
                p = en_cell.text_frame.paragraphs[0]
                p.font.color.rgb = RGBColor(255, 255, 255)
            else:
                set_cell_fill(en_cell, "F1F8F4")  # 極淺綠色
            p = en_cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.bold = is_current
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.LEFT
            en_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 中文名稱欄
            zh_cell = tbl.cell(row_idx, 2)
            zh_cell.text = zh_name
            if is_current:
                set_cell_fill(zh_cell, "4A7C59")  # 深綠色（當前分類）
                p = zh_cell.text_frame.paragraphs[0]
                p.font.color.rgb = RGBColor(255, 255, 255)
            else:
                set_cell_fill(zh_cell, "F1F8F4")  # 極淺綠色
            p = zh_cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.bold = is_current
            p.font.name = 'Microsoft JhengHei'
            p.alignment = PP_ALIGN.LEFT
            zh_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        print(f"  ✓ SASB 表格已插入（代碼：{sasb_code}，類別：{sasb_name}，當前分類：{current_category_prefix}）")

    def _generate_emission_outputs(self):
        """取得 emission 輸出（優先使用 Step 2 已生成的檔案）"""
        
        # 優先使用 Step 2 傳入的輸出資料夾
        if self.emission_output_folder and os.path.exists(self.emission_output_folder):
            print(f"  ✓ 使用 Step 2 Emission 輸出：{self.emission_output_folder}")
            
            # 尋找表格和圓餅圖
            table_files = glob.glob(os.path.join(self.emission_output_folder, "Emission_Table_*.pptx"))
            pie_files = glob.glob(os.path.join(self.emission_output_folder, "Emission_PieChart*.png"))
            
            results = {}
            if table_files:
                results["table_pptx"] = max(table_files, key=os.path.getmtime)
                print(f"  ✓ 找到表格：{os.path.basename(results['table_pptx'])}")
            if pie_files:
                results["pie_chart"] = max(pie_files, key=os.path.getmtime)
                print(f"  ✓ 找到圓餅圖：{os.path.basename(results['pie_chart'])}")
            
            if results:
                return results
        
        # 回退：呼叫引擎重新生成
        try:
            print("  ℹ 未找到 Step 2 輸出，重新生成...")
            from emission_pptx import generate_all
            results = generate_all()
            return results
        except Exception as e:
            print(f"  ⚠ Emission 引擎錯誤: {e}")
            return None

    def generate_ghg_pages(self):
        """生成溫室氣體排放管理頁面"""
        print("\n[生成溫室氣體管理頁面]")
        
        # 先生成 emission 引擎輸出
        emission_results = self._generate_emission_outputs()
        
        # Page 11: 4.5 碳盤查表（左文右表）
        section_slide = self._add_slide()
        self._add_title(section_slide, "4.5 碳盤查表")
        
        # 左邊：文字說明
        ghg_text = self.content_engine.generate_ghg_calculation_method(self.config)
        self._add_text_box(section_slide, ghg_text,
                          left=LEFT_CONTENT_LEFT, 
                          top=CONTENT_TOP, 
                          width=CONTENT_WIDTH, 
                          height=CONTENT_HEIGHT,
                          font_size=Pt(12))
        
        # 右邊：排放表格（縮減50%，靠右對齊）
        create_emission_table_on_slide_right(section_slide)
        
        # Page 12: 電力使用與節能政策（使用 emission 圓餅圖）
        electricity_text = self.content_engine.generate_electricity_policy(self.config)
        if emission_results and "pie_chart" in emission_results:
            # 使用 emission 生成的圓餅圖
            pie_chart_path = emission_results["pie_chart"]
            self._create_left_image_right_text_slide_full_path(
                "電力使用與節能政策",
                str(pie_chart_path),
                electricity_text
            )
        else:
            # 使用預設圖片
            self._create_left_text_right_image_slide(
                "電力使用與節能政策",
                electricity_text,
                ENVIRONMENT_IMAGE_MAPPING['ghg_pie']
            )
        
        # Page 13: 節能技術措施
        efficiency_text = self.content_engine.generate_energy_efficiency_measures(self.config)
        self._create_left_text_right_image_slide(
            "節能技術措施",
            efficiency_text,
            ENVIRONMENT_IMAGE_MAPPING['ghg_bar']
        )
        
        print("✓ 溫室氣體管理頁面完成")

    def generate_environmental_management_pages(self):
        """生成環境管理頁面"""
        print("\n[生成環境管理頁面]")
        
        # 4.6 綠色植育
        plant_text = self.content_engine.generate_green_planting_program(self.config)
        self._create_left_image_right_text_slide(
            "4.6 綠色植育",
            ENVIRONMENT_IMAGE_MAPPING['plant'],
            plant_text
        )
        
        # 4.7 水資源管理
        water_text = self.content_engine.generate_water_management(self.config)
        self._create_left_text_right_image_slide(
            "4.7 水資源管理",
            water_text,
            ENVIRONMENT_IMAGE_MAPPING['water']
        )
        
        # 4.8 廢棄物管理
        waste_text = self.content_engine.generate_waste_management(self.config)
        self._create_left_text_right_image_slide(
            "4.8 廢棄物管理",
            waste_text,
            ENVIRONMENT_IMAGE_MAPPING['waste']
        )
        
        # 4.9 環境教育與合作
        education_text = self.content_engine.generate_environmental_education(self.config)
        self._create_left_text_right_image_slide(
            "4.9 環境教育與合作",
            education_text,
            ENVIRONMENT_IMAGE_MAPPING['ecowork']
        )
        
        print("✓ 環境管理頁面完成")

    def generate(self):
        """生成完整環境篇 PPTX 報告"""
        print("\n" + "="*50)
        print("開始生成 ESG 環境篇 PPTX 報告")
        print(f"產業：{self.industry}")
        print("="*50)
        
        # 設定動態排放數據
        if self.emission_data:
            try:
                from assets.emission_pptx import set_emission_data
                set_emission_data(self.emission_data)
            except Exception as e:
                print(f"  ⚠ 無法載入動態排放數據: {e}")
        
        self.generate_cover_page()
        self.generate_policy_pages()
        self.generate_tcfd_pages()
        self.generate_ghg_pages()
        self.generate_environmental_management_pages()
        
        print("\n" + "="*50)
        print("PPTX 報告生成完成！")
        print(f"總共 {len(self.prs.slides)} 頁投影片")
        print("="*50)
        
        return self.prs

    def save(self, filename):
        """儲存 PPTX 檔案"""
        self.prs.save(filename)
        print(f"✓ 已儲存：{filename}")


# 測試用
if __name__ == "__main__":
    from datetime import datetime
    
    # 可以指定模板路徑
    template_path = r"C:\Users\User\Downloads\templet_resaved.pptx"
    
    engine = EnvironmentPPTXEngine(template_path=template_path)
    pptx = engine.generate()
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"output/ESG環境篇_{timestamp}.pptx"
    engine.save(filename)

