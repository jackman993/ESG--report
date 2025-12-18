"""
Component to render the ESG Core Pillars table (company page 1.3).
"""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


PILLARS = [
    {
        "pillar": "地球",
        "focus": "氣候、排放、資源韌性",
        "initiatives": [
            "淨零路線圖與階段性碳預算",
            "採用再生能源與能源效率升級",
            "水資源管理與自然正向倡議",
        ],
        "color": RGBColor(16, 185, 129),
    },
    {
        "pillar": "產品",
        "focus": "循環設計、包裝、永續採購",
        "initiatives": [
            "生命週期評估指導生態設計決策",
            "封閉循環包裝與回收內容目標",
            "供應商 ESG 資格認證與材料追溯性",
        ],
        "color": RGBColor(245, 158, 11),
    },
    {
        "pillar": "人員",
        "focus": "倫理、能力建構、包容文化",
        "initiatives": [
            "多元、公平與包容計畫",
            "綠色與數位職位的再技能培訓",
            "福祉、健康安全與社區參與",
        ],
        "color": RGBColor(30, 58, 138),
    },
]


class ESGPillarsTableComponent:
    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        left_cm: float = 18.0,
        top_cm: float = 3.0,
        width_cm: float = 13.5,
        height_cm: float = 12.5,
    ):
        rows = len(PILLARS) + 1
        cols = 3
        table = slide.shapes.add_table(
            rows,
            cols,
            Cm(left_cm),
            Cm(top_cm),
            Cm(width_cm),
            Cm(height_cm),
        ).table

        column_widths = (3.5, 4.5, width_cm - 8.0)
        for idx, width in enumerate(column_widths):
            table.columns[idx].width = Cm(width)

        headers = ["ESG 支柱", "策略重點", "關鍵倡議"]
        for idx, header in enumerate(headers):
            cell = table.cell(0, idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 41, 55)
            for para in cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10.5)
                para.font.bold = True
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.alignment = PP_ALIGN.CENTER
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        alt_bg = RGBColor(226, 232, 240)
        for row_idx, pillar in enumerate(PILLARS, start=1):
            cell_pillar = table.cell(row_idx, 0)
            cell_pillar.text = pillar["pillar"]
            cell_pillar.fill.solid()
            cell_pillar.fill.fore_color.rgb = pillar["color"]
            for para in cell_pillar.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(11)
                para.font.bold = True
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.alignment = PP_ALIGN.CENTER
                para.space_before = Pt(0)
                para.space_after = Pt(2)

            cell_focus = table.cell(row_idx, 1)
            cell_focus.text = pillar["focus"]
            cell_focus.fill.solid()
            cell_focus.fill.fore_color.rgb = alt_bg
            for para in cell_focus.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10)
                para.font.bold = False
                para.font.color.rgb = RGBColor(55, 65, 81)
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

            cell_init = table.cell(row_idx, 2)
            init_text = "\n".join(f"• {item}" for item in pillar["initiatives"])
            cell_init.text = init_text
            for para in cell_init.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(9.5)
                para.font.bold = False
                para.font.color.rgb = RGBColor(45, 55, 72)
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        return table


