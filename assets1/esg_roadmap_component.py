"""
Component to render the ESG roadmap timeline (company pages 1.5 part 1 & 2).
"""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
ROADMAP_EVENTS = [
    {
        "year": "2022",
        "title": "基礎承諾",
        "subtitle": "建立與治理",
        "items": [
            "治理：建立董事會層級 ESG 委員會與內部指導委員會。",
            "承諾：宣布營運與融資排放的淨零目標。",
            "衡量：啟動範疇一與範疇二排放基準線與資料收集。",
        ],
    },
    {
        "year": "2023",
        "title": "風險與策略整合",
        "subtitle": "政策發展與風險地圖",
        "items": [
            "風險：將氣候風險整合至企業風險管理架構。",
            "社會：強制全面倫理與反貪腐培訓。",
            "財務：發展綠色與永續金融架構與產業政策。",
            "董事會：達成 100% 董事參與氣候培訓。",
        ],
    },
    {
        "year": "2024",
        "title": "營運脫碳",
        "subtitle": "執行與新指標",
        "items": [
            "營運：啟動淨零執行計畫與能源效率升級。",
            "供應鏈：推出永續採購政策與 ESG 篩選。",
            "治理：將高階主管薪酬連結長期 ESG KPI。",
            "資料：完成融資排放（範疇三）基準線衡量。",
        ],
    },
    {
        "year": "2025",
        "title": "前瞻性合規",
        "subtitle": "精進與全球標準",
        "items": [
            "合規：對齊 IFRS S1/S2 (ISSB) 揭露要求。",
            "參與：正式化高影響力客戶的轉型規劃參與。",
            "技術：將 AI 倫理與資料治理嵌入產品開發。",
            "檢討：取得主要 ESG 目標與揭露的外部保證。",
        ],
    },
]


class ESGRoadmapComponent:
    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        phase: str = "phase1",
        left_cm: float = 18.0,
        top_cm: float = 3.0,
        width_cm: float = 13.5,
        box_height_cm: float = 5.5,
        gap_cm: float = 0.8,
    ):
        if phase == "phase2":
            events = ROADMAP_EVENTS[2:]
        elif phase == "all":
            events = ROADMAP_EVENTS
        else:
            events = ROADMAP_EVENTS[:2]

        box_width = Cm(width_cm)
        left = Cm(left_cm)
        current_top = Cm(top_cm)
        circle_radius = Cm(0.4)
        line_color = RGBColor(55, 65, 81)

        for event in events:
            # vertical line segment
            line = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                left - circle_radius / 2,
                current_top,
                Cm(0.1),
                Cm(box_height_cm),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = line_color
            line.line.fill.background()

            # circular marker
            circle = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                left - circle_radius,
                current_top + Cm(0.3),
                circle_radius * 2,
                circle_radius * 2,
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(16, 185, 129)
            circle.line.color.rgb = RGBColor(16, 185, 129)

            # content box
            box = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                left,
                current_top,
                box_width,
                Cm(box_height_cm),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(241, 245, 249)
            box.line.color.rgb = RGBColor(148, 163, 184)
            box.line.width = Pt(1.2)
            box.shadow.inherit = False

            text_frame = box.text_frame
            text_frame.clear()

            title_para = text_frame.paragraphs[0]
            title_para.text = f"{event['year']} — {event['title']}"
            title_para.font.name = self.font_name
            title_para.font.size = Pt(14)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(30, 41, 59)
            title_para.alignment = PP_ALIGN.LEFT
            title_para.space_after = Pt(2)

            subtitle_para = text_frame.add_paragraph()
            subtitle_para.text = event["subtitle"]
            subtitle_para.font.name = self.font_name
            subtitle_para.font.size = Pt(10.5)
            subtitle_para.font.bold = False
            subtitle_para.font.color.rgb = RGBColor(75, 85, 99)
            subtitle_para.alignment = PP_ALIGN.LEFT
            subtitle_para.space_before = Pt(0)
            subtitle_para.space_after = Pt(4)

            for item in event["items"]:
                clean_item = item.replace("**", "")
                bullet = text_frame.add_paragraph()
                bullet.text = clean_item
                bullet.level = 1
                bullet.font.name = self.font_name
                bullet.font.size = Pt(10)
                bullet.font.color.rgb = RGBColor(55, 65, 81)
                bullet.alignment = PP_ALIGN.LEFT
                bullet.space_before = Pt(0)
                bullet.space_after = Pt(1)

            current_top += Cm(box_height_cm + gap_cm)

