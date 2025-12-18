"""
PPT component to render the Board skill coverage table (company section page 1.4).
"""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


SKILL_ROWS = [
    {"skill": "氣候 / 永續", "current": 85, "target": 100, "notes": "TCFD 監督、轉型計畫審查"},
    {"skill": "AI / 資料倫理", "current": 60, "target": 100, "notes": "負責任 AI、資料治理、隱私"},
    {"skill": "網路安全", "current": 65, "target": 100, "notes": "韌性、事件應變、供應商安全"},
    {"skill": "人力資本管理", "current": 75, "target": 100, "notes": "接班、多元公平包容、勞動力轉型"},
    {"skill": "財務 / 風險", "current": 100, "target": 100, "notes": "資本配置、風險偏好、控制"},
]


class BoardSkillTableComponent:
    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        left_cm: float = 18.0,
        top_cm: float = 3.0,
        width_cm: float = 13.0,
        height_cm: float = 12.5,
    ):
        rows = len(SKILL_ROWS) + 1
        cols = 4
        table = slide.shapes.add_table(
            rows,
            cols,
            Cm(left_cm),
            Cm(top_cm),
            Cm(width_cm),
            Cm(height_cm),
        ).table

        column_widths = (4.0, 2.7, 2.7, width_cm - 9.4)
        for idx, w in enumerate(column_widths):
            table.columns[idx].width = Cm(w)

        headers = ["技能領域", "目前涵蓋率 (%)", "目標涵蓋率 (%)", "關鍵治理備註"]
        header_colors = [RGBColor(15, 76, 129)] * cols

        for idx, header in enumerate(headers):
            cell = table.cell(0, idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_colors[idx]
            for para in cell.text_frame.paragraphs:
                para.text = header
                para.font.name = self.font_name
                para.font.size = Pt(11)
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        for row_idx, row_data in enumerate(SKILL_ROWS, start=1):
            skill_cell = table.cell(row_idx, 0)
            skill_cell.text = row_data["skill"]
            skill_cell.fill.solid()
            skill_cell.fill.fore_color.rgb = RGBColor(224, 242, 255)
            for para in skill_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10.5)
                para.font.bold = True
                para.font.color.rgb = RGBColor(31, 41, 55)
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

            current_cell = table.cell(row_idx, 1)
            current_cell.text = f"{row_data['current']}%"
            current_cell.fill.solid()
            current_cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
            for para in current_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(11)
                para.font.bold = True
                para.font.color.rgb = RGBColor(12, 74, 110)
                para.alignment = PP_ALIGN.CENTER

            target_cell = table.cell(row_idx, 2)
            target_cell.text = f"{row_data['target']}%"
            target_cell.fill.solid()
            target_cell.fill.fore_color.rgb = RGBColor(236, 253, 245)
            for para in target_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(11)
                para.font.bold = True
                para.font.color.rgb = RGBColor(22, 101, 52)
                para.alignment = PP_ALIGN.CENTER

            notes_cell = table.cell(row_idx, 3)
            notes_cell.text = row_data["notes"]
            for para in notes_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10)
                para.font.bold = False
                para.font.color.rgb = RGBColor(55, 65, 81)
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        return table


