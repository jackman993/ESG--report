"""
Component to render the Products and Services in Response to Sustainability Goals table (company page 3.3).
"""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HEADERS = [
    "產品 / 服務類別",
    "描述",
    "應用技術",
]

ROWS = [
    {
        "category": "能源管理系統",
        "description": (
            "端對端能源績效管理，識別減排機會、優化負載曲線，並建立可信的再生能源路線圖（現場太陽能、購電協議、再生能源憑證）。"
            "解決方案整合基準線建立、目標設定和符合 ISO 50001 的 M&V 流程，同時支持範疇二減排策略（市場與位置基礎會計）。"
            "推動成本節省、削峰填谷和排放減量，提供可稽核的證據以獲得保證。"
        ),
        "technology": (
            "物聯網計量、BMS/SCADA 整合、邊緣分析、雲端資料湖、自動異常檢測、數位孿生，"
            "以及可匯出至 ESG 報告系統的儀表板 API。"
        ),
    },
    {
        "category": "智慧製造服務",
        "description": (
            "資料驅動的生產優化，減少廢料、水與能源強度，以及非計畫性停機。"
            "應用精實原則與即時品質分析，最小化材料損失並改善整體設備效率。"
            "透過副產品價值化與可回收性設計指導支持循環經濟，同時嵌入職業健康安全與倫理採購的風險控制。"
        ),
        "technology": (
            "預測性維護的機器學習、流程挖掘、MES/ERP 整合、電腦視覺品質控制，"
            "以及連接到合規知識庫的數位工作指示。"
        ),
    },
    {
        "category": "數位學習平台",
        "description": (
            "企業學習生態系統，擴展 ESG 素養、氣候風險意識，以及綠色技能再培訓與提升。"
            "實現公平的培訓機會、支持人才流動，並培養責任與創新文化。"
            "包含脫碳、TCFD/ISSB 揭露準備和供應商參與的精選課程。"
        ),
        "technology": (
            "雲端學習管理系統、適應性學習、微認證、學習成果分析、內容創作工具，"
            "以及用於跨組織協作的安全單一登入整合。"
        ),
    },
    {
        "category": "員工健康計畫",
        "description": (
            "涵蓋身體、心理和社會層面的整體健康與福祉架構，以提升生產力、留任率和心理安全。"
            "提供保密支援、早期風險檢測、人因工程介入和公平福利。"
            "計畫 KPI 對齊人力資本報告與多元公平包容目標，同時保護隱私。"
        ),
        "technology": (
            "穿戴裝置整合、隱私保護分析、遠距醫療平台、參與應用程式，"
            "以及用於匿名趨勢報告與影響衡量的安全人力資源資訊系統連結。"
        ),
    },
    {
        "category": "綠色供應鏈",
        "description": (
            "供應商賦能與績效管理計畫，推進低碳採購、廢棄物最小化和負責任採購。"
            "整合行為準則、入職、稽核和能力建構。"
            "支持產品生命週期評估、包裝減量，以及入站物流優化，以降低範疇三排放並提供追溯性。"
        ),
        "technology": (
            "供應商計分卡、追溯性平台、具 ESG 標準的電子採購、路線優化、生命週期資料庫，"
            "以及對齊溫室氣體議定書類別映射的自動計算引擎。"
        ),
    },
    {
        "category": "環境感測系統",
        "description": (
            "持續監測環境空氣品質、溫度、濕度和噪音，以指導場址層級風險控制和社區影響管理。"
            "實現極端天氣與熱壓力的早期預警，支持生物多樣性倡議，"
            "並提供符合 TCFD/ISSB 的氣候風險評估資料。"
        ),
        "technology": (
            "低功耗物聯網感測器、邊緣處理、衛星資料融合、地理空間分析、事件驅動警示，"
            "以及連接到風險儀表板和事件應變工作流程的安全資料管道。"
        ),
    },
]


class ProdSDGTableComponent:
    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        left_cm: float = 2.5,
        top_cm: float = 3.0,
        width_cm: float = 27.0,
        height_cm: float = 13.5,
    ):
        rows = len(ROWS) + 1
        cols = len(HEADERS)
        table = slide.shapes.add_table(
            rows,
            cols,
            Cm(left_cm),
            Cm(top_cm),
            Cm(width_cm),
            Cm(height_cm),
        ).table

        column_widths = (5.0, 13.34, width_cm - 5.0 - 13.34)
        for idx, width in enumerate(column_widths):
            table.columns[idx].width = Cm(width)

        header_bg = RGBColor(15, 76, 129)
        header_font = RGBColor(255, 255, 255)

        for idx, header in enumerate(HEADERS):
            cell = table.cell(0, idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg
            for para in cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(11)
                para.font.bold = True
                para.font.color.rgb = header_font
                para.alignment = PP_ALIGN.CENTER
                para.space_before = Pt(0)
                para.space_after = Pt(4)

        body_color = RGBColor(55, 65, 81)
        accent_bg = RGBColor(229, 243, 255)

        for row_idx, data in enumerate(ROWS, start=1):
            category_cell = table.cell(row_idx, 0)
            category_cell.text = data["category"]
            category_cell.fill.solid()
            category_cell.fill.fore_color.rgb = accent_bg
            for para in category_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10)
                para.font.bold = True
                para.font.color.rgb = RGBColor(31, 41, 55)
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

            description_cell = table.cell(row_idx, 1)
            description_cell.text = data["description"]
            for para in description_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(9.5)
                para.font.color.rgb = body_color
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

            technology_cell = table.cell(row_idx, 2)
            technology_cell.text = data["technology"]
            for para in technology_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(9.5)
                para.font.color.rgb = body_color
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        return table
