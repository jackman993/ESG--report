"""
Risk Management flowchart component for PPT slides (company page 3.5).
"""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR

BOXES = [
    {
        "key": "supervisory",
        "title": "監事會",
        "items": ["要求調查", "半年風險檢討"],
        "column": "left",
        "row": 0,
        "fill": (0, 153, 153),
    },
    {
        "key": "audit",
        "title": "稽核委員會",
        "items": ["確認控制有效性", "季度進度報告"],
        "column": "right",
        "row": 0,
        "fill": (0, 153, 153),
    },
    {
        "key": "board",
        "title": "管理委員會",
        "items": [],
        "column": "center",
        "row": 1,
        "fill": (0, 115, 207),
    },
    {
        "key": "cesr",
        "title": "安全與風險委員會 (CESR)",
        "items": ["設定風險偏好", "核准風險政策"],
        "column": "left",
        "row": 2,
        "fill": (0, 169, 224),
    },
    {
        "key": "disclosure",
        "title": "揭露與內部控制委員會",
        "items": ["追蹤控制有效性", "報告風險回應"],
        "column": "right",
        "row": 2,
        "fill": (0, 169, 224),
    },
    {
        "key": "owners",
        "title": "風險負責人",
        "items": ["執行緩解計畫", "報告狀態指標"],
        "column": "center",
        "row": 3,
        "fill": (27, 59, 122),
    },
]

CONNECTORS = [
    ("supervisory", "board"),
    ("audit", "board"),
    ("board", "cesr"),
    ("board", "disclosure"),
    ("cesr", "owners"),
    ("disclosure", "owners"),
]


class RiskFlowchartComponent:
    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        left_cm: float = 19.0,
        top_cm: float = 3.0,
        width_cm: float = 13.5,
        box_height_cm: float = 3.0,
        vertical_gap_cm: float = 0.6,
        box_corner_radius_cm: float = 0.3,
    ):
        layout = self._build_layout(
            left_cm,
            top_cm,
            width_cm,
            box_height_cm,
            vertical_gap_cm,
        )
        shapes = {}
        for box in BOXES:
            shapes[box["key"]] = self._draw_box(slide, box, layout, box_corner_radius_cm)
        self._draw_connectors(slide, shapes)

    def _build_layout(self, left_cm, top_cm, width_cm, box_height_cm, vertical_gap_cm):
        widths = {
            "left": min(4.0, width_cm * 0.32),
            "right": min(4.0, width_cm * 0.32),
            "center": min(5.6, width_cm * 0.46),
        }
        centers = {
            "left": left_cm + width_cm * 0.18,
            "center": left_cm + width_cm * 0.5,
            "right": left_cm + width_cm * 0.82,
        }
        rows = {}
        for box in BOXES:
            row_index = box["row"]
            if row_index not in rows:
                rows[row_index] = top_cm + row_index * (box_height_cm + vertical_gap_cm)
        return {
            "centers": centers,
            "widths": widths,
            "row_top": rows,
            "box_height": box_height_cm,
        }

    def _draw_box(self, slide, box, layout, corner_radius_cm):
        column = box["column"]
        width_cm = layout["widths"][column]
        left_cm = layout["centers"][column] - width_cm / 2
        top_cm = layout["row_top"][box["row"]]
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Cm(left_cm),
            Cm(top_cm),
            Cm(width_cm),
            Cm(layout["box_height"]),
        )
        fill_color = RGBColor(*box["fill"])
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1.5)
        max_dim = max(width_cm, layout["box_height"])
        shape.adjustments[0] = min(0.5, corner_radius_cm / max_dim)
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        title_para = frame.paragraphs[0]
        title_para.text = box["title"]
        title_para.font.bold = True
        title_para.font.size = Pt(14)
        title_para.font.name = self.font_name
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = 1
        title_para.space_after = Pt(4)
        for item in box["items"]:
            para = frame.add_paragraph()
            para.text = item
            para.font.bold = False
            para.font.size = Pt(11)
            para.font.name = self.font_name
            para.font.color.rgb = RGBColor(235, 244, 250)
            para.alignment = 1
            para.space_before = Pt(0)
            para.space_after = Pt(2)
        return shape

    def _draw_connectors(self, slide, shapes):
        for start_key, end_key in CONNECTORS:
            start_shape = shapes[start_key]
            end_shape = shapes[end_key]
            start_x = start_shape.left + start_shape.width / 2
            start_y = start_shape.top + start_shape.height
            end_x = end_shape.left + end_shape.width / 2
            end_y = end_shape.top
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
            )
            connector.line.color.rgb = RGBColor(33, 90, 146)
            connector.line.width = Pt(2.5)
            arrow_width = Cm(0.45)
            arrow_height = Cm(0.5)
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
                end_x - arrow_width / 2,
                end_y - arrow_height,
                arrow_width,
                arrow_height,
            )
            arrow.rotation = 180
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(33, 90, 146)
            arrow.line.fill.background()
