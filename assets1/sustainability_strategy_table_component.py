"""
Component to render the Sustainability Strategy & Action Plan table (company page 1.2).
"""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


HEADERS = [
    "核心永續支柱",
    "策略方向",
    "未來工作承諾",
]

ROWS = [
    {
        "pillar": "支柱一\n綠色企業",
        "direction": "綠色平台發展",
        "commitments": [
            "與策略夥伴建立永續供應鏈管理",
            "提升供應商與合作夥伴的綠色認證率",
            "開發 ESG 產品地圖，提供綠色與永續產品",
            "實施「城市到城市」永續發展計畫",
            "推出數位平台追蹤永續指標",
        ],
    },
    {
        "pillar": "支柱一\n內部管理",
        "direction": "內部管理強化",
        "commitments": [
            "取得 ISO 14001 與 ISO 50001 認證",
            "優化能源消耗並減少用水",
            "在所有設施推動廢棄物減量計畫",
            "為關鍵場所取得綠建築認證",
            "實施能源效率計畫並設定減量目標",
        ],
    },
    {
        "pillar": "支柱二\n永續營運",
        "direction": "促進長期價值",
        "commitments": [
            "將永續 KPI 嵌入產品生命週期審查",
            "擴展循環商業模式與低碳產品",
            "強化負責任採購與供應商參與",
            "擴大氣候風險評估至各資產類別",
            "將永續性整合至投資決策",
        ],
    },
    {
        "pillar": "支柱三\n創新與影響",
        "direction": "加速轉型",
        "commitments": [
            "投資綠色創新與新興技術",
            "與學術界合作孵化 ESG 解決方案",
            "支持社會影響計畫與包容性成長",
            "強化 ESG 報告的資料治理",
            "發布年度進度檢討與利害關係人回饋",
        ],
    },
]


class SustainabilityStrategyTableComponent:
    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        left_cm: float = 18.0,
        top_cm: float = 3.0,
        width_cm: float = 13.5,
        height_cm: float = 13.5,
    ):
        rows = len(ROWS) + 1
        cols = 3
        table = slide.shapes.add_table(
            rows,
            cols,
            Cm(left_cm),
            Cm(top_cm),
            Cm(width_cm),
            Cm(height_cm),
        ).table

        column_widths = (4.2, 3.8, width_cm - 8.0)
        for idx, width in enumerate(column_widths):
            table.columns[idx].width = Cm(width)

        header_bg = RGBColor(15, 76, 129)
        header_font = RGBColor(255, 255, 255)

        for idx, header in enumerate(HEADERS):
            cell = table.cell(0, idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg
            for para in cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10.5)
                para.font.bold = True
                para.font.color.rgb = header_font
                para.alignment = PP_ALIGN.CENTER
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        for row_idx, data in enumerate(ROWS, start=1):
            pillar_cell = table.cell(row_idx, 0)
            pillar_cell.text = data["pillar"]
            pillar_cell.fill.solid()
            pillar_cell.fill.fore_color.rgb = RGBColor(229, 243, 255)
            for para in pillar_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10.5)
                para.font.bold = True
                para.font.color.rgb = RGBColor(31, 41, 55)
                para.alignment = PP_ALIGN.CENTER
                para.space_before = Pt(0)
                para.space_after = Pt(2)

            direction_cell = table.cell(row_idx, 1)
            direction_cell.text = data["direction"]
            direction_cell.fill.solid()
            direction_cell.fill.fore_color.rgb = RGBColor(238, 242, 255)
            for para in direction_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10)
                para.font.bold = True
                para.font.color.rgb = RGBColor(45, 55, 72)
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

            commitments_cell = table.cell(row_idx, 2)
            commitments_text = "\n".join(f"• {item}" for item in data["commitments"])
            commitments_cell.text = commitments_text
            for para in commitments_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(9.5)
                para.font.bold = False
                para.font.color.rgb = RGBColor(55, 65, 81)
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        return table


