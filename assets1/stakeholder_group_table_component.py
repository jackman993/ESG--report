"""
Component to render the stakeholder group table (company page 2.2).
"""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

STAKEHOLDER_DATA = [
    {
        "group": "投資人與股東",
        "topics": "氣候韌性、TCFD 揭露、財務表現、治理結構、ESG 連結高階主管薪酬。",
        "channels": "股東常會、季度財報會議、ESG 評等機構參與、企業溝通。",
    },
    {
        "group": "員工與潛在人才",
        "topics": "人才吸引與留任、多元公平包容、健康安全、倫理文化與反貪腐培訓。",
        "channels": "員工脈動調查、內部培訓工作坊、全員大會、申訴程序。",
    },
    {
        "group": "客戶與顧客",
        "topics": "資料隱私與網路安全、公平對待、產品安全、服務可及性、永續金融選項。",
        "channels": "客戶回饋管道、服務專線、神秘客調查、NPS 滿意度調查。",
    },
    {
        "group": "監管機構與政府",
        "topics": "遵循 HKMA、SFC、GDPR/POPIA；風險管理整合；對齊 ISSB/ESRS。",
        "channels": "監管提交、季度合規報告、產業論壇、政策倡議。",
    },
    {
        "group": "供應商與合作夥伴",
        "topics": "負責任供應鏈管理、人權、倫理採購、永續採購合規。",
        "channels": "供應商行為準則確認、第三方風險評估、盡職調查稽核。",
    },
]


class StakeholderGroupTableComponent:
    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        left_cm: float = 3.0,
        top_cm: float = 3.0,
        width_cm: float = 27.0,
        height_cm: float = 13.5,
    ):
        rows = len(STAKEHOLDER_DATA) + 1
        cols = 3
        table = slide.shapes.add_table(
            rows,
            cols,
            Cm(left_cm),
            Cm(top_cm),
            Cm(width_cm),
            Cm(height_cm),
        ).table

        column_widths = (width_cm * 0.2, width_cm * 0.4, width_cm * 0.4)
        for idx, width in enumerate(column_widths):
            table.columns[idx].width = Cm(width)

        headers = [
            "利害關係人群體",
            "重大關切議題",
            "參與管道與頻率",
        ]

        header_bg = RGBColor(243, 244, 246)
        header_color = RGBColor(31, 41, 55)

        for idx, header in enumerate(headers):
            cell = table.cell(0, idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg
            for para in cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10.5)
                para.font.bold = True
                para.font.color.rgb = header_color
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        body_color = RGBColor(55, 65, 81)

        for row_idx, data in enumerate(STAKEHOLDER_DATA, start=1):
            group_cell = table.cell(row_idx, 0)
            group_cell.text = data["group"]
            for para in group_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10)
                para.font.bold = True
                para.font.color.rgb = header_color
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

            topics_cell = table.cell(row_idx, 1)
            topics_cell.text = data["topics"]
            for para in topics_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10)
                para.font.bold = False
                para.font.color.rgb = body_color
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

            channels_cell = table.cell(row_idx, 2)
            channels_cell.text = data["channels"]
            for para in channels_cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10)
                para.font.bold = False
                para.font.color.rgb = body_color
                para.alignment = PP_ALIGN.LEFT
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        return table

