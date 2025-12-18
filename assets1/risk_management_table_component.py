"""
Component to render the Material Issues and Risk Management table (company page 3.6).
"""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HEADERS = [
    "重大議題",
    "風險描述",
    "風險類型",
    "風險嚴重性",
    "風險可能性",
    "緩解措施",
]

ROWS = [
    {
        "issue": "永續供應鏈",
        "risk": (
            "因地緣政治不穩定或原物料短缺造成的供應中斷，面臨合規缺失與範疇三資料缺口。"
        ),
        "type": "營運",
        "severity": "中等",
        "likelihood": "高",
        "mitigation": (
            "建立供應商風險評估、多元化採購、進行定期 ESG 稽核，並以通過率 KPI 追蹤矯正行動。"
        ),
    },
    {
        "issue": "氣候策略",
        "risk": (
            "實體與轉型風險可能影響營運、聲譽和成本基礎，隨著政策、碳定價和市場期望轉變。"
        ),
        "type": "策略",
        "severity": "高",
        "likelihood": "中等",
        "mitigation": (
            "設定科學基礎目標、投資再生能源與儲能、整合 TCFD/ISSB 揭露，並執行情境分析。"
        ),
    },
    {
        "issue": "廢棄物管理",
        "risk": (
            "低效率的廢棄物處理可能損害聲譽、增加監管風險，並在缺乏循環設計的情況下侵蝕價值。"
        ),
        "type": "環境",
        "severity": "低",
        "likelihood": "中等",
        "mitigation": (
            "擴展分類、回收與副產品價值化；設定轉向目標與供應商包裝減量 KPI。"
        ),
    },
    {
        "issue": "能源效率",
        "risk": (
            "能源成本上升與清潔能源取得受限，增加營運波動性與氣候暴露。"
        ),
        "type": "營運",
        "severity": "中等",
        "likelihood": "高",
        "mitigation": (
            "升級至高效率設備、部署持續監控、削峰填谷，以及符合 ISO 50001 並具驗證 M&V 的流程。"
        ),
    },
    {
        "issue": "健康與安全",
        "risk": (
            "職場傷害或心理社會風險可能影響生產力與員工福祉。"
        ),
        "type": "營運",
        "severity": "低",
        "likelihood": "中等",
        "mitigation": (
            "進行針對性培訓、檢查、未遂事件通報、人因工程介入，以及以 TRIR KPI 為基礎的資料驅動活動。"
        ),
    },
]


class RiskManagementTableComponent:
    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        left_cm: float = 2.5,
        top_cm: float = 3.0,
        width_cm: float = 27.0,
        height_cm: float = 13.0,
    ):
        rows = len(ROWS) + 1
        cols = len(HEADERS)
        table = slide.shapes.add_table(
            rows,
            cols,
            Cm(left_cm),
            Cm(top_cm),
            Cm(width_cm),
            Cm(height_cm),
        ).table

        column_widths = (
            width_cm * 0.18,
            width_cm * 0.28,
            width_cm * 0.12,
            width_cm * 0.12,
            width_cm * 0.12,
            width_cm * 0.18,
        )
        for idx, width in enumerate(column_widths):
            table.columns[idx].width = Cm(width)

        header_bg = RGBColor(31, 78, 120)
        header_font = RGBColor(255, 255, 255)

        for idx, header in enumerate(HEADERS):
            cell = table.cell(0, idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg
            for para in cell.text_frame.paragraphs:
                para.font.name = self.font_name
                para.font.size = Pt(10.5)
                para.font.bold = True
                para.font.color.rgb = header_font
                para.alignment = PP_ALIGN.CENTER
                para.space_before = Pt(0)
                para.space_after = Pt(2)

        stripe_bg = RGBColor(242, 242, 242)
        body_color = RGBColor(55, 65, 81)

        for row_idx, data in enumerate(ROWS, start=1):
            cells = [
                data["issue"],
                data["risk"],
                data["type"],
                data["severity"],
                data["likelihood"],
                data["mitigation"],
            ]
            for col_idx, text in enumerate(cells):
                cell = table.cell(row_idx, col_idx)
                cell.text = text
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = stripe_bg
                for para in cell.text_frame.paragraphs:
                    para.font.name = self.font_name
                    para.font.size = Pt(9.5 if col_idx in (1, 5) else 9)
                    para.font.color.rgb = body_color
                    para.alignment = PP_ALIGN.LEFT if col_idx != 2 else PP_ALIGN.CENTER
                    para.space_before = Pt(0)
                    para.space_after = Pt(2)

        return table

