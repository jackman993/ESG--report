"""
Reusable table component for slide 6.5 Social Action Plan.
"""
from typing import Optional

from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


DEFAULT_ROWS = [
    {
        "initiative": "社區識字實驗室",
        "owner": "企業社會責任與教育辦公室",
        "status": "已完成",
        "next_step": "發布影響報告並將模式擴展至鄉村學校。",
    },
    {
        "initiative": "員工志工團",
        "owner": "人力資源與文化",
        "status": "進行中",
        "next_step": "擴展專業導師庫並增加週末班次。",
    },
    {
        "initiative": "鄰里健康篩檢",
        "owner": "公共衛生合作夥伴團隊",
        "status": "進度落後",
        "next_step": "確保額外護理人員，更新溝通計畫。",
    },
    {
        "initiative": "循環廢棄物聯盟",
        "owner": "永續辦公室",
        "status": "設計中",
        "next_step": "確認試點場址和供應商認證標準。",
    },
]


class SocialActionPlanTable:
    """Render a four-column action tracker on the slide."""

    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        left_cm: float = 2.3,
        top_cm: float = 3.0,
        width_cm: float = 14.0,
        height_cm: Optional[float] = None,
        rows=None,
    ):
        rows = rows or DEFAULT_ROWS
        height_cm = height_cm or 13.0

        table = slide.shapes.add_table(
            rows=len(rows) + 1,
            cols=4,
            left=Cm(left_cm),
            top=Cm(top_cm),
            width=Cm(width_cm),
            height=Cm(height_cm),
        ).table

        column_widths = (4.0, 3.5, 3.0, 3.5)
        for idx, width in enumerate(column_widths):
            table.columns[idx].width = Cm(width)

        headers = ("倡議", "負責單位", "狀態", "下一步")
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = self.font_name
                paragraph.font.size = Pt(11.5)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(40, 40, 40)
                paragraph.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(222, 224, 228)

        status_colors = {
            "已完成": RGBColor(82, 139, 94),
            "進行中": RGBColor(63, 120, 177),
            "進度落後": RGBColor(191, 120, 43),
            "設計中": RGBColor(120, 120, 120),
            # 保留英文狀態以向後兼容
            "Completed": RGBColor(82, 139, 94),
            "On Track": RGBColor(63, 120, 177),
            "Lagging": RGBColor(191, 120, 43),
            "In Design": RGBColor(120, 120, 120),
        }

        for row_idx, row in enumerate(rows, start=1):
            values = (
                row["initiative"],
                row["owner"],
                row["status"],
                row["next_step"],
            )
            for col_idx, value in enumerate(values):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = self.font_name
                    paragraph.font.size = Pt(10.5)
                    paragraph.font.color.rgb = RGBColor(66, 66, 66)
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)
                    paragraph.alignment = PP_ALIGN.LEFT

                if col_idx == 2:  # status column styling
                    fill_color = status_colors.get(value, RGBColor(120, 120, 120))
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = fill_color
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.alignment = PP_ALIGN.CENTER
                elif row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(246, 246, 246)

        return table

