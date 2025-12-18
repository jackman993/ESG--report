"""
Reusable table component for slide 6.1 Community Engagement & Social Investment.
"""
from typing import Optional

from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


DEFAULT_ROWS = [
    {
        "program": "未來技能",
        "focus": "為青年提供 STEM 和數位素養學院、學徒制，以及與認證合作夥伴共同提供的中途轉職再培訓。",
        "iso": "ISO 26000 6.8.3 — 教育與文化",
    },
    {
        "program": "社區健康",
        "focus": "與在地健康網絡共同提供行動診所、心理健康活動和預防性篩檢服務。",
        "iso": "ISO 26000 6.8.5 — 社會投資",
    },
    {
        "program": "環境",
        "focus": "流域復育、氣候韌性補助，以及解決社區廢棄物挑戰的循環經濟試點計畫。",
        "iso": "ISO 26000 6.5 & 6.8.5 — 環境與社會投資",
    },
    {
        "program": "動物福利",
        "focus": "獸醫外展服務、人道收容所，以及強化生物多樣性管理的野生動物共存倡議。",
        "iso": "ISO 26000 6.8.5 — 社會投資",
    },
]


class CommunityInvestmentTable:
    """Draw a three-column programme overview table on the slide."""

    def __init__(self, prs, font_name: str = "Calibri"):
        self.prs = prs
        self.font_name = font_name

    def add_to_slide(
        self,
        slide,
        left_cm: float = 18.0,
        top_cm: float = 3.0,
        width_cm: float = 14.0,
        height_cm: Optional[float] = None,
        rows=None,
    ):
        rows = rows or DEFAULT_ROWS
        height_cm = height_cm or 13.0

        table = slide.shapes.add_table(
            rows=len(rows) + 1,
            cols=3,
            left=Cm(left_cm),
            top=Cm(top_cm),
            width=Cm(width_cm),
            height=Cm(height_cm),
        ).table

        column_widths = (4.0, 6.5, 3.5)
        for idx, width in enumerate(column_widths):
            table.columns[idx].width = Cm(width)

        headers = ("計畫", "影響重點", "ISO 對齊")
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = self.font_name
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(45, 45, 45)
                paragraph.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 230, 230)

        for row_idx, row in enumerate(rows, start=1):
            programme = row["program"]
            focus = row["focus"]
            iso = row["iso"]
            for col_idx, value in enumerate((programme, focus, iso)):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = self.font_name
                    paragraph.font.size = Pt(10.5)
                    paragraph.font.color.rgb = RGBColor(66, 66, 66)
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)
                    paragraph.alignment = PP_ALIGN.LEFT
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

        return table

